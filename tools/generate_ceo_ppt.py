"""
generate_ceo_ppt.py
Single combined CEO presentation — Appraisal Data + AL sheet (last slide).
Run: source .venv/bin/activate && python tools/generate_ceo_ppt.py
"""
import openpyxl
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime

BASE = "/Applications/XAMPP/xamppfiles/htdocs/minerva"
SRC  = f"{BASE}/Book1 (003).xlsx"
OUT  = f"{BASE}/tools/CEO_Appraisal_Presentation.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0D, 0x1B, 0x3E)
BLUE    = RGBColor(0x1A, 0x6B, 0xB5)
LBLUE   = RGBColor(0x5B, 0x9B, 0xD5)
GOLD    = RGBColor(0xF0, 0xA8, 0x00)
GREEN   = RGBColor(0x2E, 0x8B, 0x57)
LGREEN  = RGBColor(0x70, 0xAD, 0x47)
ORANGE  = RGBColor(0xE0, 0x6C, 0x00)
RED     = RGBColor(0xC0, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHT  = RGBColor(0xF5, 0xF7, 0xFA)
LGREY   = RGBColor(0xEB, 0xED, 0xF0)
MGREY   = RGBColor(0xB0, 0xB8, 0xC8)
DGREY   = RGBColor(0x33, 0x3C, 0x4E)

# ── Data loading ──────────────────────────────────────────────────────────────
def load(sheet_name):
    wb = openpyxl.load_workbook(SRC, read_only=True)
    rows = list(wb[sheet_name].iter_rows(values_only=True))
    hdrs = rows[0]
    return [dict(zip(hdrs, r)) for r in rows[1:] if any(v for v in r)]

def g(row, key):
    for k in [key, key.strip(), key + ' ']:
        if k in row and row[k] is not None:
            return str(row[k]).strip()
    return ''

def stats(data):
    total = len(data)
    ratings  = Counter(g(r,'Rating') for r in data)
    promos   = Counter(g(r,'Promotion') for r in data)
    appr     = Counter(g(r,'Apprisal') for r in data)
    bands    = Counter(r.get('Band') for r in data)
    directors= Counter(g(r,'Director') for r in data if g(r,'Director'))
    clients  = Counter(g(r,'Client') for r in data if g(r,'Client'))
    desigs   = Counter(g(r,'Actual Designation') for r in data if g(r,'Actual Designation'))
    supers   = Counter(g(r,'Supervisor') for r in data if g(r,'Supervisor'))

    dir_r = defaultdict(Counter)
    dir_p = defaultdict(Counter)
    band_r= defaultdict(Counter)
    band_p= defaultdict(Counter)
    desig_p = defaultdict(Counter)
    for row in data:
        d = g(row,'Director')
        dir_r[d][g(row,'Rating')] += 1
        dir_p[d][g(row,'Promotion')] += 1
        band_r[row.get('Band')][g(row,'Rating')] += 1
        band_p[row.get('Band')][g(row,'Promotion')] += 1
        desig_p[g(row,'Actual Designation')][g(row,'Promotion')] += 1

    return dict(
        total=total, ratings=ratings, promos=promos, appr=appr,
        bands=bands, directors=directors, clients=clients, desigs=desigs,
        supers=supers, dir_r=dir_r, dir_p=dir_p,
        band_r=band_r, band_p=band_p, desig_p=desig_p,
        r4=ratings.get('4',0), r3=ratings.get('3',0), r2=ratings.get('2',0),
        na=ratings.get('NA',0),
        yes=promos.get('Yes',0), no=promos.get('No',0),
        full=appr.get('Fully Eligible',0),
        part=appr.get('Partially Eligible',0),
        noel=appr.get('Not Eligible',0),
        rated=ratings.get('4',0)+ratings.get('3',0)+ratings.get('2',0),
    )

# ── PPT helpers ───────────────────────────────────────────────────────────────
def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def box(slide, l, t, w, h, fill, line_c=None, line_w=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_c:
        s.line.color.rgb = line_c
        if line_w: s.line.width = Pt(line_w)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, sz=11, bold=False, color=DGREY,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = wrap
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def section_label(slide, label, l=0.35, t=1.15):
    txt(slide, label.upper(), l, t, 12, 0.3, 7.5, True, MGREY)

def divider(slide, l, t, w, color=GOLD):
    box(slide, l, t, w, 0.025, color)

def header(slide, title, subtitle="", bg_c=NAVY):
    box(slide, 0, 0, 13.33, 0.95, bg_c)
    box(slide, 0, 0.95, 13.33, 0.03, GOLD)
    txt(slide, title, 0.35, 0.10, 11, 0.52, 20, True, WHITE)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.63, 11, 0.28, 9, False, MGREY, italic=True)

def kpi(slide, label, value, note, l, t, w=1.85, h=1.3, fill=BLUE, note_c=None):
    box(slide, l, t, w, h, fill, line_c=WHITE, line_w=0.3)
    txt(slide, label, l+0.08, t+0.08, w-0.16, 0.32, 8, True, GOLD, PP_ALIGN.CENTER)
    txt(slide, str(value), l+0.08, t+0.38, w-0.16, 0.48, 26, True, WHITE, PP_ALIGN.CENTER)
    txt(slide, note, l+0.08, t+h-0.36, w-0.16, 0.32, 7.5, False,
        note_c if note_c else RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.CENTER)

def col_chart(slide, title, cats, series, l, t, w, h, colors=None, legend=True):
    if colors is None: colors = [BLUE, LGREEN, RGBColor(0xBF,0xBF,0xBF), GOLD, ORANGE]
    cd = ChartData()
    cd.categories = [str(c) for c in cats]
    for name, vals in series:
        cd.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    ).chart
    chart.has_title  = bool(title)
    chart.has_legend = legend
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(9)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
    return chart

def bar_chart(slide, title, cats, vals, l, t, w, h, color=BLUE):
    cd = ChartData()
    cd.categories = [str(c) for c in cats]
    cd.add_series('', vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    ).chart
    chart.has_title  = bool(title)
    chart.has_legend = False
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(9)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    for ser in chart.series:
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = color
    return chart

def pie_chart(slide, cats, vals, l, t, w, h, colors=None):
    cd = ChartData()
    cd.categories = [str(c) for c in cats]
    cd.add_series('', vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(l), Inches(t), Inches(w), Inches(h), cd
    ).chart
    chart.has_title  = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return chart

def mini_table(slide, headers, rows, col_ws, l, t, row_h=0.3):
    """Draw a lightweight table without borders — header + data rows."""
    # Header
    x = l
    for h_txt, cw in zip(headers, col_ws):
        box(slide, x, t, cw, row_h, NAVY)
        txt(slide, h_txt, x+0.05, t+0.05, cw-0.1, row_h-0.05, 7.5, True, WHITE, PP_ALIGN.CENTER)
        x += cw
    # Rows
    for ri, row_vals in enumerate(rows):
        top = t + row_h*(ri+1)
        fill_c = LGREY if ri%2==0 else WHITE
        x = l
        for val, cw in zip(row_vals, col_ws):
            box(slide, x, top, cw, row_h, fill_c, RGBColor(0xD0,0xD5,0xDE), 0.2)
            al = PP_ALIGN.LEFT if x==l else PP_ALIGN.CENTER
            txt(slide, str(val), x+0.05, top+0.04, cw-0.1, row_h-0.04, 8, False, DGREY, al)
            x += cw

def footer(slide, note=""):
    box(slide, 0, 7.38, 13.33, 0.02, BLUE)
    src = note if note else f"Source: Book1 (003).xlsx  |  Confidential  |  {datetime.now().strftime('%B %Y')}"
    txt(slide, src, 0.25, 7.42, 12.8, 0.22, 7, False, MGREY)

# ─────────────────────────────────────────────────────────────────────────────
# Load & compute
# ─────────────────────────────────────────────────────────────────────────────
print("Loading data…")
ad   = load('Apprisal Data')
ald  = load('AL')
A    = stats(ad)
AL   = stats(ald)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, NAVY)
box(sl, 0, 0, 0.08, 7.5, GOLD)           # left accent bar
box(sl, 0.08, 5.78, 13.25, 0.04, GOLD)   # lower divider

txt(sl, "ANNUAL APPRAISAL CYCLE", 1.0, 1.4, 11, 0.55, 13, True, GOLD)
txt(sl, "Performance Review — FY 2024–25", 1.0, 2.05, 11, 0.9, 36, True, WHITE)
txt(sl, "Comprehensive analysis of appraisal outcomes across all employees,\n"
        "covering ratings, eligibility, promotions, and director performance.",
    1.0, 3.1, 10, 0.75, 12, False, MGREY, italic=True)
txt(sl, f"Prepared for the CEO & Leadership Council  |  {datetime.now().strftime('%B %Y')}",
    1.0, 6.05, 11, 0.35, 9, False, MGREY)
txt(sl, f"Total Employees Covered: {A['total']}   |   Data: Book1 (003).xlsx",
    1.0, 6.42, 11, 0.3, 8.5, False, RGBColor(0x60,0x70,0x90))

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Snapshot (KPI dashboard)
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Executive Snapshot", "At-a-glance KPIs for the full appraisal cycle")

# Row 1 — 4 KPI cards
gap = 0.22
kw  = (13.33 - 5*gap) / 4
kpis_r1 = [
    ("Total Employees",     A['total'],               "Organisation headcount",    NAVY),
    ("Fully Eligible",      f"{A['full']}  ({A['full']*100//A['total']}%)",  "Full year covered",   GREEN),
    ("Partially Eligible",  f"{A['part']}  ({A['part']*100//A['total']}%)",  "Part year covered",   ORANGE),
    ("Not Eligible",        f"{A['noel']}  ({A['noel']*100//A['total']}%)",  "Joined after Oct 2025", RED),
]
for i, (lbl, val, note, col) in enumerate(kpis_r1):
    kpi(sl, lbl, val, note, gap + i*(kw+gap), 1.15, kw, 1.35, col)

# Row 2 — 4 KPI cards
kpis_r2 = [
    ("Rated Employees",  f"{A['rated']}  ({A['rated']*100//A['total']}%)",  f"of {A['total']} total",   BLUE),
    ("Rating 4 — Exceeds", f"{A['r4']}  ({A['r4']*100//A['total']}%)",        "Top performers",        LGREEN),
    ("Rating 3 — Meets",   f"{A['r3']}  ({A['r3']*100//A['total']}%)",        "On-track performers",   LBLUE),
    ("Promoted",           f"{A['yes']}  ({A['yes']*100//A['total']}%)",       "Promotions this cycle", RGBColor(0x5B,0x5E,0xA8)),
]
for i, (lbl, val, note, col) in enumerate(kpis_r2):
    kpi(sl, lbl, val, note, gap + i*(kw+gap), 2.75, kw, 1.35, col)

# Band summary bar
band_txt = (f"Band Composition:   Band 1 = {A['bands'].get(1,0)} emp  |  "
            f"Band 2 = {A['bands'].get(2,0)} emp  |  "
            f"Band 3 = {A['bands'].get(3,0)} emp")
box(sl, gap, 4.35, 13.33-2*gap, 0.52, NAVY)
txt(sl, band_txt, gap+0.25, 4.48, 13.33-2*gap-0.3, 0.3, 9.5, False, WHITE)

# Director headcount
dirs = A['directors'].most_common()
dir_txt = "Director Portfolio:   " + "   |   ".join(f"{d} = {c} emp" for d,c in dirs)
box(sl, gap, 4.95, 13.33-2*gap, 0.52, RGBColor(0x1A, 0x2E, 0x55))
txt(sl, dir_txt, gap+0.25, 5.08, 13.33-2*gap-0.3, 0.3, 9.5, False, WHITE)

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Appraisal Eligibility & Rating
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Eligibility & Performance Ratings", "Who was rated and how they performed")

# Left: Eligibility pie
pie_chart(sl,
    ['Fully Eligible', 'Partially Eligible', 'Not Eligible'],
    [A['full'], A['part'], A['noel']],
    0.3, 1.1, 4.5, 4.2)
txt(sl, "Appraisal Eligibility", 0.3, 1.08, 4.5, 0.28, 9, True, NAVY, PP_ALIGN.CENTER)
divider(sl, 0.3, 1.35, 4.5)

# Eligibility figures below pie
elig_rows = [
    ("Fully Eligible",     A['full'],  f"{A['full']*100//A['total']}%",  "Full year in service"),
    ("Partially Eligible", A['part'],  f"{A['part']*100//A['total']}%",  "Apr–Oct 2025 joiners"),
    ("Not Eligible (NA)",  A['noel'],  f"{A['noel']*100//A['total']}%",  "Joined after Oct 2025"),
]
y = 5.42
for lbl, cnt, pct, note in elig_rows:
    txt(sl, f"● {lbl}", 0.35, y, 2.8, 0.28, 8.5, False, DGREY)
    txt(sl, f"{cnt}  ({pct})", 3.1, y, 1.7, 0.28, 8.5, True, NAVY, PP_ALIGN.RIGHT)
    y += 0.3

# Middle divider
box(sl, 5.05, 1.05, 0.02, 5.6, LGREY)

# Right: Rating column chart
col_chart(sl, "", ['Rating 4\n(Exceeds)', 'Rating 3\n(Meets)', 'Rating 2\n(Below)', 'Not Applic.\n(NA)'],
    [("Employees", [A['r4'], A['r3'], A['r2'], A['na']])],
    5.2, 1.1, 7.8, 3.6, [BLUE, LGREEN, RED, RGBColor(0xBF,0xBF,0xBF)], legend=False)
txt(sl, "Rating Distribution", 5.2, 1.08, 7.8, 0.28, 9, True, NAVY)
divider(sl, 5.2, 1.35, 7.8)

# Rating stat pills
pill_data = [
    (f"Rating 4 — Exceeds", A['r4'], f"{A['r4']*100//A['total']}%", LGREEN),
    (f"Rating 3 — Meets",   A['r3'], f"{A['r3']*100//A['total']}%", BLUE),
    (f"Rating 2 — Below",   A['r2'], f"{A['r2']*100//A['total']}%", RED),
    (f"Not Applic. (NA)",   A['na'], f"{A['na']*100//A['total']}%",  MGREY),
]
x = 5.2
for lbl, cnt, pct, col in pill_data:
    box(sl, x, 4.85, 1.82, 0.72, col)
    txt(sl, lbl,      x+0.07, 4.88, 1.7, 0.25, 7.5, True,  WHITE, PP_ALIGN.CENTER)
    txt(sl, f"{cnt}  ({pct})", x+0.07, 5.13, 1.7, 0.34, 13, True, WHITE, PP_ALIGN.CENTER)
    x += 1.95

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Director Performance
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Director Performance", "Rating and promotion outcomes by Director")

# Collect directors ordered by headcount
dir_order = [d for d,_ in A['directors'].most_common() if d]
short_names = [d.split()[-1] for d in dir_order]

# Left: Rating column chart grouped
r4v  = [A['dir_r'][d].get('4',0)  for d in dir_order]
r3v  = [A['dir_r'][d].get('3',0)  for d in dir_order]
nav  = [A['dir_r'][d].get('NA',0) for d in dir_order]
col_chart(sl, "Rating by Director", short_names,
    [("Rating 4", r4v), ("Rating 3", r3v), ("Not Eligible", nav)],
    0.3, 1.1, 6.0, 4.3, [LGREEN, BLUE, RGBColor(0xBF,0xBF,0xBF)])

# Right: Detailed director scorecard table
col_ws = [2.6, 0.75, 0.65, 0.65, 0.65, 0.85, 0.85]
tbl_hdrs = ['Director', 'Total', 'R4', 'R3', 'NA', 'Promoted', 'Promo%']
tbl_rows = []
for d in dir_order:
    dr = A['dir_r'][d]; dp = A['dir_p'][d]
    tot = sum(dr.values())
    pct = f"{dp.get('Yes',0)*100//tot}%" if tot else "0%"
    tbl_rows.append([d, tot, dr.get('4',0), dr.get('3',0), dr.get('NA',0), dp.get('Yes',0), pct])
mini_table(sl, tbl_hdrs, tbl_rows, col_ws, 6.5, 1.15, 0.38)

# Band performance table below
txt(sl, "BY BAND", 6.5, 3.1, 6.8, 0.28, 8.5, True, NAVY)
divider(sl, 6.5, 3.37, 6.8)
band_col_ws = [1.1, 0.75, 0.65, 0.65, 0.65, 0.85]
band_tbl_hdrs = ['Band','Total','R4','R3','NA','Promoted']
band_tbl_rows = []
for b in [1,2,3]:
    br = A['band_r'].get(b,Counter()); bp = A['band_p'].get(b,Counter())
    band_tbl_rows.append([f'Band {b}', A['bands'].get(b,0),
                          br.get('4',0), br.get('3',0), br.get('NA',0), bp.get('Yes',0)])
mini_table(sl, band_tbl_hdrs, band_tbl_rows, band_col_ws, 6.5, 3.45, 0.34)

# Insight box
insight = (f"Sheriff manages {A['directors'].get('Muzzafar Sheriff',0)} of {A['total']} employees "
           f"({A['directors'].get('Muzzafar Sheriff',0)*100//A['total']}% of workforce) "
           f"with {A['dir_p']['Muzzafar Sheriff'].get('Yes',0)} promotions this cycle.")
box(sl, 0.3, 5.6, 12.73, 0.62, RGBColor(0xE8, 0xF0, 0xFA), RGBColor(0x2E,0x75,0xB6), 0.5)
txt(sl, f"📌  {insight}", 0.5, 5.7, 12.5, 0.42, 9, False, NAVY)
footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Promotions Deep-Dive
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Promotions Analysis", f"  {A['yes']} employees promoted  |  {A['yes']*100//A['total']}% of total headcount")

# Left top: Pie — Promoted vs Not
pie_chart(sl, ['Promoted','Not Promoted'], [A['yes'], A['no']], 0.3, 1.1, 3.8, 3.5)
txt(sl, "Promoted vs Not Promoted", 0.3, 1.08, 3.8, 0.28, 9, True, NAVY, PP_ALIGN.CENTER)
divider(sl, 0.3, 1.35, 3.8)

# Left bottom: By Band chart
bp_yes = [A['band_p'].get(b,Counter()).get('Yes',0) for b in [1,2,3]]
bp_no  = [A['band_p'].get(b,Counter()).get('No',0)  for b in [1,2,3]]
col_chart(sl, "Promotions by Band", ['Band 1','Band 2','Band 3'],
    [("Promoted", bp_yes), ("Not Promoted", bp_no)],
    0.3, 4.75, 3.8, 2.4, [LGREEN, RGBColor(0xBF,0xBF,0xBF)])

# Centre divider
box(sl, 4.3, 1.05, 0.02, 6.1, LGREY)

# Right: Top promoted designations horizontal bars
txt(sl, "Top Promoted Designations", 4.5, 1.08, 8.5, 0.28, 9, True, NAVY)
divider(sl, 4.5, 1.35, 8.6)

top_promo = sorted([(d,c.get('Yes',0)) for d,c in A['desig_p'].items() if c.get('Yes',0)>0],
                   key=lambda x:-x[1])[:10]
max_p = top_promo[0][1] if top_promo else 1
for i, (desig, cnt) in enumerate(top_promo):
    y = 1.5 + i*0.57
    bw = cnt/max_p * 7.8
    # Background bar
    box(sl, 4.5, y+0.025, 8.6, 0.42, LGREY)
    # Value bar
    box(sl, 4.5, y+0.025, bw, 0.42, BLUE if i>0 else NAVY)
    txt(sl, desig, 4.6, y+0.04, 5.5, 0.34, 8.5, False, WHITE if bw>2 else DGREY)
    txt(sl, str(cnt), 13.0, y+0.04, 0.4, 0.34, 9, True, DGREY, PP_ALIGN.RIGHT)

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Client & Grade Overview
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Client & Grade Distribution", "Workforce composition across clients and grade levels")

# Left: Client distribution bar chart
top_cli = A['clients'].most_common(8)
bar_chart(sl, "Headcount by Client",
    [c for c,_ in top_cli], [v for _,v in top_cli],
    0.3, 1.1, 5.5, 5.5, BLUE)

# Right: Grade table + designation snapshot
txt(sl, "Grade Breakdown", 6.1, 1.08, 7.0, 0.28, 9, True, NAVY)
divider(sl, 6.1, 1.35, 7.0)

grade_col_ws = [1.4, 0.9, 0.7, 1.0]
grade_hdrs = ['Grade','Count','%','Promoted']
gp = defaultdict(int)
for row in ad:
    if g(row,'Promotion')=='Yes':
        gp[str(row.get('Grade','')).strip()] += 1
from collections import Counter as _C
grades_counter = _C(str(r.get('Grade','')).strip() for r in ad if r.get('Grade'))
grade_rows_real = []
for grade, cnt in grades_counter.most_common(12):
    pct = f"{cnt*100//A['total']}%"
    grade_rows_real.append([grade, cnt, pct, gp.get(grade,0)])
mini_table(sl, grade_hdrs, grade_rows_real, grade_col_ws, 6.1, 1.45, 0.3)

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Key Observations (all-company)
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, OFFWHT)
header(sl, "Key Observations & Recommendations", "Strategic takeaways for leadership action")

observations = [
    (NAVY,   "01", "Eligibility Gap",
     f"{A['noel']} employees ({A['noel']*100//A['total']}%) joined after Oct 2025 and are outside this "
     f"appraisal cycle. Align future hiring plans with appraisal cutoff dates."),

    (BLUE,   "02", "Rating Concentration",
     f"{A['r3']} employees rated '3 — Meets' vs {A['r4']} rated '4 — Exceeds'. "
     f"A calibration session is recommended to sharpen differentiation between performers."),

    (GREEN,  "03", "Promotion Rate",
     f"Only {A['yes']*100//A['total']}% ({A['yes']}/{A['total']}) were promoted. Band 1 has the lowest "
     f"promo rate; consider reviewing grade-progression criteria for Client Partner roles."),

    (ORANGE, "04", "Director Portfolio Concentration",
     f"Muzzafar Sheriff manages {A['directors'].get('Muzzafar Sheriff',0)} employees "
     f"({A['directors'].get('Muzzafar Sheriff',0)*100//A['total']}% of headcount). "
     f"Delegation and succession planning should be reviewed."),

    (RED,    "05", "Partial Eligibility Action",
     f"{A['part']} employees are partially eligible. Confirm pro-rated merit increments "
     f"are correctly computed and communicated before salary revision."),
]

y = 1.12
for col, num, title, body in observations:
    box(sl, 0.3, y, 0.5, 0.95, col)
    txt(sl, num, 0.3, y+0.25, 0.5, 0.4, 12, True, WHITE, PP_ALIGN.CENTER)
    box(sl, 0.82, y, 12.2, 0.95, WHITE, RGBColor(0xD0,0xD5,0xDE), 0.3)
    txt(sl, title, 0.95, y+0.07, 11.9, 0.3, 10.5, True, col)
    txt(sl, body,  0.95, y+0.38, 11.9, 0.48, 9, False, DGREY)
    y += 1.08

footer(sl)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ASSET LIVING LLC (AL Sheet) — Final Slide
# ═════════════════════════════════════════════════════════════════════════════
sl = new_slide(prs)
fill_bg(sl, RGBColor(0x07, 0x14, 0x30))   # deep navy for contrast

# Decorative side stripe
box(sl, 0, 0, 0.09, 7.5, GOLD)

# Header
box(sl, 0.09, 0, 13.24, 1.1, RGBColor(0x0D, 0x1B, 0x3E))
box(sl, 0.09, 1.1, 13.24, 0.04, GOLD)
txt(sl, "ASSET LIVING, LLC", 0.4, 0.08, 10, 0.45, 22, True, GOLD)
txt(sl, "Client Performance Spotlight  |  505 Dedicated Employees",
    0.4, 0.6, 10, 0.3, 9.5, False, MGREY, italic=True)
txt(sl, "AL Sheet", 11.8, 0.3, 1.3, 0.35, 9, False, RGBColor(0x60,0x70,0x90), PP_ALIGN.RIGHT)

# ── 6 KPI cards ───────────────────────────────────────────────────────────────
card_defs = [
    ("Total\nHeadcount",    str(AL['total']),                                   NAVY),
    ("Fully\nEligible",     f"{AL['full']}\n({AL['full']*100//AL['total']}%)",   GREEN),
    ("Partially\nEligible", f"{AL['part']}\n({AL['part']*100//AL['total']}%)",   ORANGE),
    ("Not\nEligible",       f"{AL['noel']}\n({AL['noel']*100//AL['total']}%)",   RED),
    ("Rating 4\n(Top)",     f"{AL['r4']}\n({AL['r4']*100//AL['total']}%)",       BLUE),
    ("Promoted",            f"{AL['yes']}\n({AL['yes']*100//AL['total']}%)",     LBLUE),
]
n = len(card_defs);  cw = 13.24/n
for i, (lbl, val, col) in enumerate(card_defs):
    cx = 0.09 + i*cw
    box(sl, cx+0.06, 1.22, cw-0.12, 1.3, col)
    txt(sl, lbl,         cx+0.1, 1.26, cw-0.18, 0.38, 7.5, True,  GOLD,  PP_ALIGN.CENTER)
    lines = val.split('\n')
    txt(sl, lines[0], cx+0.1, 1.62, cw-0.18, 0.45, 22, True, WHITE, PP_ALIGN.CENTER)
    if len(lines)>1:
        txt(sl, lines[1], cx+0.1, 2.05, cw-0.18, 0.3,  8,  False, RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.CENTER)

# ── Director column chart ─────────────────────────────────────────────────────
al_dirs = [d for d,_ in AL['directors'].most_common() if d]
al_short = [d.split()[-1] for d in al_dirs]
al_r4  = [AL['dir_r'][d].get('4',0)  for d in al_dirs]
al_r3  = [AL['dir_r'][d].get('3',0)  for d in al_dirs]
al_na  = [AL['dir_r'][d].get('NA',0) for d in al_dirs]
col_chart(sl, "Rating by Director", al_short,
    [("Rating 4", al_r4), ("Rating 3", al_r3), ("Not Eligible", al_na)],
    0.15, 2.65, 4.6, 4.05, [LGREEN, BLUE, RGBColor(0x60,0x60,0x60)])

# ── Band breakdown table ──────────────────────────────────────────────────────
txt(sl, "BAND BREAKDOWN", 5.0, 2.65, 5, 0.3, 8, True, GOLD)
b_hdrs = ['Band','Total','R4','R3','NA','Promo','%']
b_cws  = [0.85, 0.75, 0.6, 0.6, 0.6, 0.8, 0.75]
b_rows = []
for b in [1,2,3]:
    br = AL['band_r'].get(b,Counter()); bp = AL['band_p'].get(b,Counter())
    tot_b = AL['bands'].get(b,0)
    pct = f"{bp.get('Yes',0)*100//tot_b}%" if tot_b else "0%"
    b_rows.append([f'Band {b}', tot_b, br.get('4',0), br.get('3',0), br.get('NA',0), bp.get('Yes',0), pct])
mini_table(sl, b_hdrs, b_rows, b_cws, 5.0, 3.0, 0.32)

# ── Top promoted designations ─────────────────────────────────────────────────
txt(sl, "PROMOTED DESIGNATIONS (TOP 7)", 5.0, 4.2, 8.2, 0.3, 8, True, GOLD)
divider(sl, 5.0, 4.5, 8.22, GOLD)
al_top = sorted([(d,c.get('Yes',0)) for d,c in AL['desig_p'].items() if c.get('Yes',0)>0],
                key=lambda x:-x[1])[:7]
al_max = al_top[0][1] if al_top else 1
for i, (desig, cnt) in enumerate(al_top):
    y2 = 4.58 + i*0.4
    bw2 = cnt/al_max * 7.8
    box(sl, 5.0, y2, 8.22, 0.32, RGBColor(0x18,0x2A,0x52))
    box(sl, 5.0, y2, bw2,  0.32, BLUE if i>0 else LBLUE)
    txt(sl, f"{desig}  —  {cnt}", 5.1, y2+0.06, 7.8, 0.24, 8, False, WHITE)

# Footer
box(sl, 0, 7.38, 13.33, 0.02, GOLD)
txt(sl, f"Source: Book1 (003).xlsx — Sheet: AL  |  Asset Living LLC  |  505 records  |  Confidential  |  {datetime.now().strftime('%B %Y')}",
    0.25, 7.42, 12.8, 0.22, 7, False, RGBColor(0x50,0x60,0x80))

# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"✅  Saved → {OUT}")
print(f"   Slides: 8  (Title · KPI Snapshot · Eligibility & Rating · Director Performance · "
      f"Promotions · Client & Grade · Observations · Asset Living LLC)")
