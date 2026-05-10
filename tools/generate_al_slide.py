"""
generate_al_slide.py
Produces a single crisp slide for Asset Living, LLC — CEO view.
Run: source .venv/bin/activate && python tools/generate_al_slide.py
"""
import openpyxl
from collections import Counter, defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from datetime import datetime

BASE = "/Applications/XAMPP/xamppfiles/htdocs/minerva"
SRC  = f"{BASE}/Book1 (003).xlsx"
OUT  = f"{BASE}/tools/AssetLiving_CEO_Slide.pptx"

# ── colours ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1F, 0x35, 0x64)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
GOLD   = RGBColor(0xFF, 0xC0, 0x00)
GREEN  = RGBColor(0x70, 0xAD, 0x47)
ORANGE = RGBColor(0xFF, 0x76, 0x00)
RED    = RGBColor(0xC0, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
DGREY  = RGBColor(0x40, 0x40, 0x40)

# ── load AL sheet ─────────────────────────────────────────────────────────────
wb = openpyxl.load_workbook(SRC, read_only=True)
ws = wb['AL']
rows = list(ws.iter_rows(values_only=True))
headers = rows[0]
data = [dict(zip(headers, r)) for r in rows[1:] if any(v for v in r)]

def g(row, key):
    for k in [key, key.strip(), key + ' ']:
        if k in row and row[k] is not None:
            return str(row[k]).strip()
    return ''

total   = len(data)
ratings = Counter(g(r,'Rating') for r in data)
promos  = Counter(g(r,'Promotion') for r in data)
appr    = Counter(g(r,'Apprisal') for r in data)
bands_c = Counter(r.get('Band') for r in data)

r4   = ratings.get('4', 0)
r3   = ratings.get('3', 0)
na   = ratings.get('NA', 0)
yes  = promos.get('Yes', 0)
no   = promos.get('No', 0)
full = appr.get('Fully Eligible', 0)
part = appr.get('Partially Eligible', 0)
noel = appr.get('Not Eligible', 0)

dir_r = defaultdict(Counter)
for row in data:
    dir_r[g(row,'Director')][g(row,'Rating')] += 1

dir_p = defaultdict(Counter)
for row in data:
    dir_p[g(row,'Director')][g(row,'Promotion')] += 1

top_desig_promo = Counter()
for row in data:
    if g(row,'Promotion') == 'Yes':
        top_desig_promo[g(row,'Actual Designation')] += 1

directors = [d for d in dir_r if d]

# ── build PPT ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])

# Background
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)

# ── Header bar ────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if line is None else setattr(s.line.color, 'rgb', line)
    return s

def tb(slide, text, l, t, w, h, sz=11, bold=False, color=DGREY, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

rect(slide, 0, 0, 13.33, 1.05, NAVY)
rect(slide, 0, 1.05, 13.33, 0.04, GOLD)

tb(slide, "ASSET LIVING, LLC", 0.3, 0.08, 8, 0.45, 22, True, WHITE)
tb(slide, "Performance Appraisal Overview  |  CEO Report", 0.3, 0.55, 9, 0.38, 10, False, RGBColor(0xBD,0xD7,0xEE))
tb(slide, datetime.now().strftime("%B %Y"), 11.5, 0.35, 1.6, 0.35, 9, False, RGBColor(0x9D,0xC3,0xE6), PP_ALIGN.RIGHT)

# ── KPI strip (6 cards) ───────────────────────────────────────────────────────
card_defs = [
    ("Total\nHeadcount",  str(total),   NAVY),
    ("Fully\nEligible",   f"{full}\n{full*100//total}%",  GREEN),
    ("Part.\nEligible",   f"{part}\n{part*100//total}%",  ORANGE),
    ("Not\nEligible",     f"{noel}\n{noel*100//total}%",  RED),
    ("Rating 4\n(Top)",   f"{r4}\n{r4*100//total}%",      BLUE),
    ("Promoted",          f"{yes}\n{yes*100//total}%",    RGBColor(0x5B,0x9B,0xD5)),
]
card_w = 13.33 / len(card_defs)
for i, (lbl, val, col) in enumerate(card_defs):
    cx = i * card_w
    rect(slide, cx + 0.04, 1.18, card_w - 0.08, 1.22, col)
    # label
    tb(slide, lbl, cx + 0.08, 1.22, card_w - 0.12, 0.38, 7.5, True, GOLD, PP_ALIGN.CENTER)
    # value
    lines = val.split('\n')
    tb(slide, lines[0], cx + 0.08, 1.58, card_w - 0.12, 0.38, 20, True, WHITE, PP_ALIGN.CENTER)
    if len(lines) > 1:
        tb(slide, lines[1], cx + 0.08, 1.92, card_w - 0.12, 0.22, 8.5, False, RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.CENTER)

# ── Section: Director Performance (column chart) ──────────────────────────────
tb(slide, "DIRECTOR PERFORMANCE", 0.25, 2.62, 5, 0.28, 8.5, True, NAVY)
rect(slide, 0.25, 2.90, 4.7, 0.025, GOLD)

cd = ChartData()
cd.categories = [d.split()[-1] for d in directors]   # last name only for brevity
cd.add_series('Rating 4',  [dir_r[d].get('4',0)  for d in directors])
cd.add_series('Rating 3',  [dir_r[d].get('3',0)  for d in directors])
cd.add_series('NA',        [dir_r[d].get('NA',0) for d in directors])
chart_sp = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.2), Inches(2.95), Inches(4.8), Inches(4.1), cd
)
chart = chart_sp.chart
chart.has_legend = True
chart.has_title  = False
from pptx.enum.chart import XL_LEGEND_POSITION
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
for i, (ser, col) in enumerate(zip(chart.series, [GREEN, BLUE, RGBColor(0xBF,0xBF,0xBF)])):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = col

# ── Section: Band split table ─────────────────────────────────────────────────
tb(slide, "BAND BREAKDOWN", 5.25, 2.62, 4, 0.28, 8.5, True, NAVY)
rect(slide, 5.25, 2.90, 7.8, 0.025, GOLD)

band_headers = ['Band', 'Total', 'R4', 'R3', 'NA', 'Promoted', 'Promo %']
col_ws = [0.9, 0.85, 0.7, 0.7, 0.7, 1.0, 0.9]
band_data_rows = []
for b in [1, 2, 3]:
    from collections import Counter as C2
    br = Counter(g(r,'Rating') for r in data if r.get('Band')==b)
    bp = Counter(g(r,'Promotion') for r in data if r.get('Band')==b)
    tot_b = bands_c.get(b,0)
    pct = f"{bp.get('Yes',0)*100//tot_b}%" if tot_b else "0%"
    band_data_rows.append([f'Band {b}', tot_b, br.get('4',0), br.get('3',0), br.get('NA',0), bp.get('Yes',0), pct])

row_h = 0.32
hdr_top = 3.0
left0 = 5.3
rect(slide, left0, hdr_top, sum(col_ws), row_h, NAVY)
x = left0
for h, cw in zip(band_headers, col_ws):
    tb(slide, h, x+0.03, hdr_top+0.05, cw-0.06, row_h-0.05, 8, True, WHITE, PP_ALIGN.CENTER)
    x += cw
for ri, row_vals in enumerate(band_data_rows):
    top = hdr_top + row_h*(ri+1)
    fill_c = LGREY if ri%2==0 else WHITE
    rect(slide, left0, top, sum(col_ws), row_h, fill_c)
    x = left0
    for val, cw in zip(row_vals, col_ws):
        tb(slide, str(val), x+0.03, top+0.05, cw-0.06, row_h-0.05, 8.5, False, DGREY, PP_ALIGN.CENTER)
        x += cw

# ── Section: Top Promoted Designations ───────────────────────────────────────
tb(slide, "TOP PROMOTED DESIGNATIONS", 5.25, 4.15, 5, 0.28, 8.5, True, NAVY)
rect(slide, 5.25, 4.43, 7.8, 0.025, GOLD)

top10 = top_desig_promo.most_common(8)
for i, (desig, cnt) in enumerate(top10):
    bar_w = cnt / (top10[0][1] if top10 else 1) * 5.8
    top_r = 4.55 + i * 0.35
    rect(slide, 5.3, top_r + 0.04, bar_w, 0.24, BLUE)
    tb(slide, f"{desig}  ({cnt})", 5.35, top_r + 0.05, 7.5, 0.24, 8, False, WHITE)

# ── Section: Director promo table ─────────────────────────────────────────────
tb(slide, "DIRECTOR SUMMARY", 0.25, 7.12, 4.7, 0.28, 8, True, NAVY)
# One-liner per director
dir_line_y = 7.3
for d in directors:
    dr = dir_r[d]
    dp = dir_p[d]
    tot_d = sum(dr.values())
    short = d.replace('Muzzafar ', 'M.').replace('Payal ', 'P.').replace('Mohit ', 'Mo.')
    line = f"{short}: {tot_d} emp | R4={dr.get('4',0)}  R3={dr.get('3',0)}  NA={dr.get('NA',0)} | Promoted={dp.get('Yes',0)}"
    tb(slide, line, 0.25, dir_line_y, 6.5, 0.25, 7.5, False, DGREY)
    dir_line_y += 0.24

# ── Footer ────────────────────────────────────────────────────────────────────
rect(slide, 0, 7.35, 13.33, 0.02, BLUE)
tb(slide, f"Source: Book1 (003).xlsx — AL Sheet  |  {total} records  |  Confidential  |  {datetime.now().strftime('%B %Y')}",
   0.2, 7.38, 12, 0.22, 7, False, RGBColor(0x80,0x80,0x80))

prs.save(OUT)
print(f"✅  Saved: {OUT}")
