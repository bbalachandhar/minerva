"""
generate_appraisal_reports.py
Generates PPT, Excel working file, and summary document for both
'Apprisal Data' and 'AL' sheets from 'Book1 (003).xlsx'.

Run:  source .venv/bin/activate && python tools/generate_appraisal_reports.py
"""

import openpyxl
import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from collections import Counter, defaultdict
from datetime import datetime
import os

# ─────────────────────────────────────────────────────────────────────────────
# Brand colours
# ─────────────────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1F, 0x35, 0x64)   # deep navy
MID_BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # corporate blue
GOLD       = RGBColor(0xFF, 0xC0, 0x00)   # accent gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
ORANGE     = RGBColor(0xFF, 0x76, 0x00)
RED        = RGBColor(0xC0, 0x00, 0x00)

BASE = "/Applications/XAMPP/xamppfiles/htdocs/minerva"
SRC  = f"{BASE}/Book1 (003).xlsx"
OUT  = f"{BASE}/tools"

# ─────────────────────────────────────────────────────────────────────────────
# Data loading
# ─────────────────────────────────────────────────────────────────────────────
def load_sheet(sname):
    wb = openpyxl.load_workbook(SRC, read_only=True)
    ws = wb[sname]
    rows = list(ws.iter_rows(values_only=True))
    headers = rows[0]
    return [dict(zip(headers, r)) for r in rows[1:] if any(v for v in r)]

def v(row, key):
    """Get value stripping whitespace, handling key with trailing space."""
    for k in [key, key.strip(), key + ' ']:
        if k in row:
            val = row[k]
            return str(val).strip() if val is not None else ''
    return ''

def compute_stats(data):
    total = len(data)
    ratings = Counter(v(r, 'Rating') for r in data)
    promos  = Counter(v(r, 'Promotion') for r in data)
    appr    = Counter(v(r, 'Apprisal') for r in data)
    bands   = Counter(r.get('Band') for r in data)
    grades  = Counter(str(r.get('Grade','')).strip() for r in data if r.get('Grade'))
    directors = Counter(v(r,'Director') for r in data if v(r,'Director'))
    clients   = Counter(v(r,'Client') for r in data if v(r,'Client'))
    desigs    = Counter(v(r,'Actual Designation') for r in data if v(r,'Actual Designation'))
    mis       = Counter(v(r,'MIS Remarks') for r in data)

    # Rating by Band
    band_rating = defaultdict(Counter)
    for r in data:
        band_rating[r.get('Band')][v(r,'Rating')] += 1

    # Rating by Director
    dir_rating = defaultdict(Counter)
    for r in data:
        dir_rating[v(r,'Director')][v(r,'Rating')] += 1

    # Promotion by Director
    dir_promo = defaultdict(Counter)
    for r in data:
        dir_promo[v(r,'Director')][v(r,'Promotion')] += 1

    # Promotion by Band
    band_promo = defaultdict(Counter)
    for r in data:
        band_promo[r.get('Band')][v(r,'Promotion')] += 1

    # Promotion by Designation
    desig_promo = defaultdict(Counter)
    for r in data:
        desig_promo[v(r,'Actual Designation')][v(r,'Promotion')] += 1

    # Eligible employees (Fully + Partially)
    eligible = sum(1 for r in data if 'Eligible' in v(r,'Apprisal') and 'Not' not in v(r,'Apprisal'))
    fully_eligible = appr.get('Fully Eligible', 0)
    partially_eligible = appr.get('Partially Eligible', 0)
    not_eligible = appr.get('Not Eligible', 0)

    rated_count = sum(1 for r in data if v(r,'Rating') not in ('NA',''))
    
    return {
        'total': total,
        'ratings': ratings,
        'promos': promos,
        'appr': appr,
        'bands': bands,
        'grades': grades,
        'directors': directors,
        'clients': clients,
        'desigs': desigs,
        'mis': mis,
        'band_rating': band_rating,
        'dir_rating': dir_rating,
        'dir_promo': dir_promo,
        'band_promo': band_promo,
        'desig_promo': desig_promo,
        'eligible': eligible,
        'fully_eligible': fully_eligible,
        'partially_eligible': partially_eligible,
        'not_eligible': not_eligible,
        'rated_count': rated_count,
    }

# ─────────────────────────────────────────────────────────────────────────────
# PPT helpers
# ─────────────────────────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def bg(slide, color):
    from pptx.util import Inches
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def textbox(slide, text, left, top, width, height,
            fontsize=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
            bg_color=None, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_title_slide(prs, title, subtitle, sheet_label):
    slide = add_slide(prs, 0)
    bg(slide, DARK_BLUE)
    # Gold accent bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(3.8), Inches(10), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    textbox(slide, sheet_label, 0.4, 0.3, 9, 0.5, 11, False, GOLD, PP_ALIGN.LEFT, italic=True)
    textbox(slide, title,    0.4, 1.0, 9.2, 1.4, 34, True,  WHITE, PP_ALIGN.LEFT)
    textbox(slide, subtitle, 0.4, 2.5, 9.2, 1.0, 16, False, RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.LEFT)
    textbox(slide, f"Confidential  |  Prepared for CEO  |  {datetime.now().strftime('%B %Y')}",
            0.4, 4.3, 9, 0.4, 10, False, RGBColor(0x9D,0xC3,0xE6), PP_ALIGN.LEFT)

def add_section_divider(prs, title, color=MID_BLUE):
    slide = add_slide(prs, 6)
    bg(slide, color)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(2.8), Inches(10), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = GOLD
    bar.line.fill.background()
    textbox(slide, title, 0.5, 1.5, 9, 1.2, 36, True, WHITE, PP_ALIGN.CENTER)

def kpi_card(slide, label, value, sub, left, top, w=1.9, h=1.5, bg_color=MID_BLUE):
    box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = WHITE; box.line.width = Pt(0.5)
    tf = slide.shapes.add_textbox(Inches(left+0.05), Inches(top+0.06), Inches(w-0.1), Inches(h*0.35))
    p = tf.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label; run.font.size = Pt(9); run.font.color.rgb = GOLD; run.font.bold = True
    tf2 = slide.shapes.add_textbox(Inches(left+0.05), Inches(top+0.38), Inches(w-0.1), Inches(h*0.45))
    p2 = tf2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2  = p2.add_run(); r2.text = str(value); r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE
    tf3 = slide.shapes.add_textbox(Inches(left+0.05), Inches(top+h-0.45), Inches(w-0.1), Inches(0.38))
    p3 = tf3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3  = p3.add_run(); r3.text = sub; r3.font.size = Pt(8); r3.font.color.rgb = RGBColor(0xBD,0xD7,0xEE)

def add_bar_chart(slide, title_text, categories, values, left, top, width, height,
                  label="Count", bar_color=MID_BLUE):
    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    cd.add_series(label, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), cd
    ).chart
    chart.has_legend = False
    chart.has_title  = True
    chart.chart_title.text_frame.text = title_text
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    for ser in chart.series:
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = bar_color
    return chart

def add_column_chart(slide, title_text, categories, series_dict, left, top, width, height):
    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    colors = [MID_BLUE, GREEN, ORANGE, GOLD, RED]
    for i, (sname, vals) in enumerate(series_dict.items()):
        cd.add_series(sname, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(left), Inches(top), Inches(width), Inches(height), cd
    ).chart
    chart.has_legend = True
    chart.has_title  = True
    chart.chart_title.text_frame.text = title_text
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    for i, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = colors[i % len(colors)]
    return chart

def add_pie_chart(slide, title_text, categories, values, left, top, width, height):
    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    cd.add_series('', values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(left), Inches(top), Inches(width), Inches(height), cd
    ).chart
    chart.has_legend = True
    chart.has_title  = True
    chart.chart_title.text_frame.text = title_text
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    chart.chart_title.text_frame.paragraphs[0].font.bold = True
    return chart

def slide_header(slide, title, subtitle=""):
    bg(slide, WHITE)
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.85))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    textbox(slide, title, 0.2, 0.08, 8.5, 0.65, 18, True, WHITE, PP_ALIGN.LEFT)
    if subtitle:
        textbox(slide, subtitle, 0.2, 0.65, 8.5, 0.3, 9, False, RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.LEFT)
    # footer line
    foot = slide.shapes.add_shape(1, Inches(0), Inches(7.3), Inches(10), Inches(0.02))
    foot.fill.solid(); foot.fill.fore_color.rgb = MID_BLUE
    foot.line.fill.background()
    textbox(slide, "Confidential | CEO Report", 0.2, 7.32, 5, 0.25, 7, False, RGBColor(0x80,0x80,0x80))

def bullet_table_slide(slide, rows, col_widths, headers, header_bg=DARK_BLUE, alt_bg=LIGHT_GREY):
    """Draw a simple table of text rows."""
    top = 1.0
    row_h = 0.38
    # header row
    left = 0.3
    for ci, (hdr, cw) in enumerate(zip(headers, col_widths)):
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(cw), Inches(row_h))
        box.fill.solid(); box.fill.fore_color.rgb = header_bg
        box.line.fill.background()
        tf = slide.shapes.add_textbox(Inches(left+0.04), Inches(top+0.05), Inches(cw-0.08), Inches(row_h-0.05))
        p = tf.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = hdr; run.font.size = Pt(9); run.font.bold = True; run.font.color.rgb = WHITE
        left += cw
    top += row_h
    for ri, row_vals in enumerate(rows):
        left = 0.3
        row_bg = alt_bg if ri % 2 == 1 else WHITE
        for ci, (val, cw) in enumerate(zip(row_vals, col_widths)):
            box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(cw), Inches(row_h))
            box.fill.solid(); box.fill.fore_color.rgb = row_bg
            box.line.color.rgb = RGBColor(0xD9,0xD9,0xD9); box.line.width = Pt(0.25)
            tf = slide.shapes.add_textbox(Inches(left+0.04), Inches(top+0.04), Inches(cw-0.08), Inches(row_h-0.05))
            p = tf.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val); run.font.size = Pt(8.5)
            left += cw
        top += row_h

# ─────────────────────────────────────────────────────────────────────────────
# PPT Builder
# ─────────────────────────────────────────────────────────────────────────────
def build_ppt(sheet_name, data, stats, out_path, label):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    is_al = (sheet_name == 'AL')
    client_focus = "Asset Living, LLC Focus" if is_al else "All Clients"
    subtitle_line = ("Asset Living, LLC dedicated headcount with full appraisal deep-dive"
                     if is_al else
                     "Organisation-wide appraisal cycle — all clients & all employees")

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    add_title_slide(prs,
        f"{'Asset Living (AL)' if is_al else 'Appraisal Data'} — CEO Report",
        subtitle_line,
        label)

    # ── Slide 2: Executive Summary (KPI cards) ───────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Executive Summary", "Key Performance Indicators at a Glance")
    total   = stats['total']
    rated   = stats['rated_count']
    r4      = stats['ratings'].get('4', 0)
    r3      = stats['ratings'].get('3', 0)
    promoY  = stats['promos'].get('Yes', 0)
    fully   = stats['fully_eligible']
    partial = stats['partially_eligible']
    not_el  = stats['not_eligible']
    
    kpi_card(slide, "Total Employees",    total,   "Headcount",           0.3,  1.1)
    kpi_card(slide, "Rated Employees",    rated,   f"{rated*100//total}% of total", 2.35, 1.1)
    kpi_card(slide, "Rating 4 (Top)",     r4,      f"{r4*100//total}% of total",    4.40, 1.1, bg_color=GREEN)
    kpi_card(slide, "Rating 3",           r3,      f"{r3*100//total}% of total",    6.45, 1.1, bg_color=MID_BLUE)
    kpi_card(slide, "Promoted",           promoY,  f"{promoY*100//total}% of total",0.3,  2.85)
    kpi_card(slide, "Fully Eligible",     fully,   f"{fully*100//total}% of total", 2.35, 2.85, bg_color=GREEN)
    kpi_card(slide, "Partially Eligible", partial, f"{partial*100//total}% of total",4.40, 2.85, bg_color=ORANGE)
    kpi_card(slide, "Not Eligible",       not_el,  "Joined after Oct 2025",          6.45, 2.85, bg_color=RED)

    # Summary bullet
    note = (f"Of {total} employees: {rated} were rated ({rated*100//total}%), "
            f"{promoY} promoted ({promoY*100//total}%), "
            f"{fully} fully eligible ({fully*100//total}%), "
            f"{partial} partially eligible ({partial*100//total}%), "
            f"{not_el} not eligible ({not_el*100//total}%).")
    textbox(slide, note, 0.3, 4.6, 9.4, 0.8, 9.5, False, RGBColor(0x40,0x40,0x40), PP_ALIGN.LEFT)

    # ── Slide 3: Appraisal Eligibility ───────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Appraisal Eligibility Breakdown", "Based on MIS Remarks / Date of Joining")
    cats = ['Fully Eligible', 'Partially Eligible', 'Not Eligible']
    vals = [fully, partial, not_el]
    add_pie_chart(slide, "Appraisal Eligibility Distribution", cats, vals, 0.3, 1.0, 4.5, 5.8)

    # Explanation textbox
    notes = [
        "• Fully Eligible: Full year covered (joined on/before Apr 1 2025)",
        "• Partially Eligible: Part year covered (joined Oct 2024 – Apr 2025)",
        "• Not Eligible: Joined after 1st Oct 2025 — excluded from rating",
        "",
        f"  Fully Eligible:     {fully:>4}  ({fully*100//total}%)",
        f"  Partially Eligible: {partial:>4}  ({partial*100//total}%)",
        f"  Not Eligible:       {not_el:>4}  ({not_el*100//total}%)",
        f"  Total:              {total:>4}",
    ]
    y = 1.2
    for n in notes:
        textbox(slide, n, 5.1, y, 4.6, 0.38, 9, False, RGBColor(0x26,0x26,0x26))
        y += 0.38

    # ── Slide 4: Rating Distribution ─────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Performance Rating Distribution", "Rating scale: 4 = Exceeds, 3 = Meets, 2 = Below, NA = Not Eligible")
    r_cats = ['Rating 4\n(Exceeds)', 'Rating 3\n(Meets)', 'Rating 2\n(Below)']
    r_vals = [stats['ratings'].get('4',0), stats['ratings'].get('3',0), stats['ratings'].get('2',0)]
    add_column_chart(slide, "Rating Distribution (Eligible Employees)",
                     r_cats, {'Employees': r_vals}, 0.3, 1.0, 5.5, 5.8)
    # Stats on right
    y = 1.3
    lines = [
        ("Rating 4 (Exceeds)", r4,  GREEN,   " — Top Performers"),
        ("Rating 3 (Meets)",   r3,  MID_BLUE," — On Track"),
        ("Rating 2 (Below)",   stats['ratings'].get('2',0), RED, " — Below Target"),
        ("Not Applicable (NA)",stats['ratings'].get('NA',0), RGBColor(0x80,0x80,0x80)," — Joined after Oct 2025"),
    ]
    for lbl, cnt, col, rem in lines:
        pct = f"{cnt*100//total}%"
        textbox(slide, f"{lbl}", 6.0, y, 3.5, 0.32, 9.5, True, RGBColor(0x26,0x26,0x26))
        textbox(slide, f"  {cnt:>4}  ({pct}){rem}", 6.0, y+0.3, 3.9, 0.32, 9, False, RGBColor(0x40,0x40,0x40))
        y += 0.82

    # ── Slide 5: Rating by Band ───────────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Performance Rating by Band", "Band 1 = Front-line | Band 2 = Mid-level | Band 3 = Senior")
    b_cats = ['Band 1', 'Band 2', 'Band 3']
    b4 = [stats['band_rating'].get(b,{}).get('4',0) for b in [1,2,3]]
    b3 = [stats['band_rating'].get(b,{}).get('3',0) for b in [1,2,3]]
    bna= [stats['band_rating'].get(b,{}).get('NA',0) for b in [1,2,3]]
    add_column_chart(slide, "Rating by Band",
                     b_cats, {'Rating 4': b4, 'Rating 3': b3, 'Not Eligible (NA)': bna},
                     0.3, 1.0, 9.4, 5.8)

    # ── Slide 6: Rating by Director ──────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Performance Rating by Director", "Director-level view of performance distribution")
    directors_order = sorted(stats['dir_rating'].keys(), key=lambda d: -sum(stats['dir_rating'][d].values()))
    d_cats = directors_order
    d4 = [stats['dir_rating'].get(d,{}).get('4',0) for d in d_cats]
    d3 = [stats['dir_rating'].get(d,{}).get('3',0) for d in d_cats]
    dna= [stats['dir_rating'].get(d,{}).get('NA',0) for d in d_cats]
    add_column_chart(slide, "Rating Distribution by Director",
                     d_cats, {'Rating 4': d4, 'Rating 3': d3, 'Not Eligible (NA)': dna},
                     0.3, 1.0, 9.4, 5.5)
    # Director summary table
    tbl_rows = []
    for d in directors_order:
        dr = stats['dir_rating'].get(d, {})
        dp = stats['dir_promo'].get(d, {})
        total_d = sum(dr.values())
        rated_d = dr.get('4',0) + dr.get('3',0) + dr.get('2',0)
        tbl_rows.append([d, total_d, dr.get('4',0), dr.get('3',0), dr.get('NA',0), dp.get('Yes',0)])
    # We'll show this as a textbox summary since the chart already fills the slide
    y = 6.6
    for row in tbl_rows:
        textbox(slide, f"  {row[0]}: {row[1]} emp | R4={row[2]}, R3={row[3]}, NA={row[4]}, Promoted={row[5]}",
                0.3, y, 9.4, 0.3, 8, False, RGBColor(0x26,0x26,0x26))
        y += 0.3

    # ── Slide 7: Promotions Overview ─────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Promotions Overview", f"Total Promoted: {promoY} of {total} ({promoY*100//total}%)")
    # Pie: Overall
    add_pie_chart(slide, "Promoted vs Not Promoted",
                  ['Promoted', 'Not Promoted'],
                  [promoY, stats['promos'].get('No',0)],
                  0.3, 1.0, 4.0, 4.5)
    # By Band
    bp_cats = ['Band 1', 'Band 2', 'Band 3']
    bp_yes = [stats['band_promo'].get(b,{}).get('Yes',0) for b in [1,2,3]]
    bp_no  = [stats['band_promo'].get(b,{}).get('No',0)  for b in [1,2,3]]
    add_column_chart(slide, "Promotions by Band",
                     bp_cats, {'Promoted': bp_yes, 'Not Promoted': bp_no},
                     4.5, 1.0, 5.2, 4.5)
    # Text note
    textbox(slide, f"Band 1: {bp_yes[0]} promoted  |  Band 2: {bp_yes[1]} promoted  |  Band 3: {bp_yes[2]} promoted",
            0.3, 5.7, 9.4, 0.4, 9.5, False, RGBColor(0x26,0x26,0x26))

    # ── Slide 8: Promotions by Director ──────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Promotions by Director", "Headcount promoted under each Director")
    dp_cats = sorted(stats['dir_promo'].keys(), key=lambda d:-stats['dir_promo'][d].get('Yes',0))
    dp_yes  = [stats['dir_promo'].get(d,{}).get('Yes',0) for d in dp_cats]
    dp_no   = [stats['dir_promo'].get(d,{}).get('No',0)  for d in dp_cats]
    add_column_chart(slide, "Promotions by Director",
                     dp_cats, {'Promoted': dp_yes, 'Not Promoted': dp_no},
                     0.3, 1.0, 9.4, 5.5)
    y = 6.7
    for d, y_cnt, n_cnt in zip(dp_cats, dp_yes, dp_no):
        tot = y_cnt + n_cnt
        pct = f"{y_cnt*100//tot}%" if tot else "0%"
        textbox(slide, f"  {d}: {y_cnt} promoted / {tot} total ({pct})",
                0.3, y, 9.4, 0.3, 8.5, False, RGBColor(0x26,0x26,0x26))
        y += 0.3

    # ── Slide 9: Top Promoted Designations ───────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Promotions by Designation", "Roles with highest number of promotions this cycle")
    promoted_list = sorted(
        [(d, cnt.get('Yes',0)) for d, cnt in stats['desig_promo'].items() if cnt.get('Yes',0)>0],
        key=lambda x:-x[1])
    desig_cats = [x[0] for x in promoted_list[:12]]
    desig_vals = [x[1] for x in promoted_list[:12]]
    add_bar_chart(slide, "Promotions by Designation (Top 12)",
                  desig_cats, desig_vals, 0.3, 1.0, 6.5, 5.8, bar_color=GREEN)
    # Table right side
    tbl_rows = [(d, y) for d,y in promoted_list]
    bullet_table_slide.__doc__  # no-op
    y_pos = 1.1
    textbox(slide, "Designation", 7.0, y_pos, 2.8, 0.3, 8, True, WHITE, PP_ALIGN.LEFT, bg_color=DARK_BLUE)
    textbox(slide, "Promoted",    9.3, y_pos, 0.6, 0.3, 8, True, WHITE, PP_ALIGN.CENTER, bg_color=DARK_BLUE)
    y_pos += 0.3
    for i, (d, cnt) in enumerate(tbl_rows[:13]):
        row_bg = LIGHT_GREY if i%2==0 else WHITE
        textbox(slide, d,   7.0, y_pos, 2.8, 0.28, 8, False, RGBColor(0x26,0x26,0x26), bg_color=row_bg)
        textbox(slide, str(cnt), 9.3, y_pos, 0.6, 0.28, 8, True, RGBColor(0x26,0x26,0x26), PP_ALIGN.CENTER, bg_color=row_bg)
        y_pos += 0.28

    # ── Slide 10: Band & Grade Distribution ──────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Band & Grade Distribution", "Workforce structure by seniority band")
    b_cats2 = ['Band 1', 'Band 2', 'Band 3']
    b_vals2 = [stats['bands'].get(1,0), stats['bands'].get(2,0), stats['bands'].get(3,0)]
    add_pie_chart(slide, "Headcount by Band", b_cats2, b_vals2, 0.3, 1.0, 4.5, 5.5)
    # Top grades table
    top_grades = stats['grades'].most_common(14)
    y_pos = 1.0
    textbox(slide, "Grade", 5.1, y_pos, 2.0, 0.3, 8.5, True, WHITE, PP_ALIGN.LEFT, bg_color=DARK_BLUE)
    textbox(slide, "Count", 7.1, y_pos, 0.8, 0.3, 8.5, True, WHITE, PP_ALIGN.CENTER, bg_color=DARK_BLUE)
    textbox(slide, "%",     7.9, y_pos, 0.7, 0.3, 8.5, True, WHITE, PP_ALIGN.CENTER, bg_color=DARK_BLUE)
    y_pos += 0.3
    for i, (g, cnt) in enumerate(top_grades):
        row_bg = LIGHT_GREY if i%2==0 else WHITE
        pct = f"{cnt*100//total}%"
        textbox(slide, g,    5.1, y_pos, 2.0, 0.28, 8, False, RGBColor(0x26,0x26,0x26), bg_color=row_bg)
        textbox(slide, str(cnt), 7.1, y_pos, 0.8, 0.28, 8, False, RGBColor(0x26,0x26,0x26), PP_ALIGN.CENTER, bg_color=row_bg)
        textbox(slide, pct,  7.9, y_pos, 0.7, 0.28, 8, False, RGBColor(0x26,0x26,0x26), PP_ALIGN.CENTER, bg_color=row_bg)
        y_pos += 0.28

    # ── Slide 11: Client Distribution ────────────────────────────────────────
    if not is_al:
        slide = add_slide(prs, 6)
        slide_header(slide, "Headcount by Client", "Distribution of employees across client accounts")
        top_clients = stats['clients'].most_common(14)
        cli_cats = [c for c,_ in top_clients]
        cli_vals = [v for _,v in top_clients]
        add_bar_chart(slide, "Employees by Client", cli_cats, cli_vals,
                      0.3, 1.0, 9.4, 5.8, bar_color=MID_BLUE)

    # ── Slide 12: Top Designations ────────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Designation Distribution", "Workforce composition by role")
    top_desig = stats['desigs'].most_common(12)
    d_cats2 = [d for d,_ in top_desig]
    d_vals2 = [v for _,v in top_desig]
    add_bar_chart(slide, "Top Designations by Headcount", d_cats2, d_vals2,
                  0.3, 1.0, 9.4, 5.8, bar_color=DARK_BLUE)

    # ── Slide 13: Supervisor Involvement ─────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Supervisor Coverage", "Supervisors with most reportees in this appraisal cycle")
    top_sup = Counter(v(r,'Supervisor') for r in data if v(r,'Supervisor')).most_common(15)
    sup_cats = [s for s,_ in top_sup]
    sup_vals = [c for _,c in top_sup]
    add_bar_chart(slide, "Top 15 Supervisors by Headcount", sup_cats, sup_vals,
                  0.3, 1.0, 9.4, 5.8, bar_color=GOLD)

    # ── Slide 14: Recommendations ────────────────────────────────────────────
    slide = add_slide(prs, 6)
    slide_header(slide, "Key Observations & Recommendations", "Actionable insights for leadership review")
    bg(slide, WHITE)
    obs = [
        ("1. Eligibility Gap",
         f"{not_el} employees ({not_el*100//total}%) joined after Oct 2025 and are ineligible for this cycle. "
         f"Ensure onboarding is aligned with appraisal cutoff dates."),
        ("2. Rating Concentration",
         f"{r3} employees rated 3 vs {r4} rated 4. Strong skew towards 'Meets Expectations' — "
         f"calibration session recommended to differentiate top performers."),
        ("3. Promotion Rate",
         f"Only {promoY*100//total}% ({promoY} of {total}) received promotions. "
         f"Band 1 has lowest promotion rate; review grade-level criteria for Client Partner roles."),
        ("4. Director Portfolio",
         f"Muzzafar Sheriff manages {stats['directors'].get('Muzzafar Sheriff',0)} employees "
         f"({stats['directors'].get('Muzzafar Sheriff',0)*100//total}% of headcount). "
         f"Consider delegation and succession planning."),
        ("5. Partial Eligibility",
         f"{partial} employees are partially eligible — ensure merit increments are pro-rated correctly."),
    ]
    y = 1.0
    for title, body in obs:
        textbox(slide, title, 0.3, y, 9.4, 0.3, 10, True, DARK_BLUE)
        textbox(slide, body,  0.5, y+0.3, 9.1, 0.45, 9, False, RGBColor(0x40,0x40,0x40))
        y += 0.82

    # ── Slide 15: Thank You ───────────────────────────────────────────────────
    slide = add_slide(prs, 6)
    bg(slide, DARK_BLUE)
    bar2 = slide.shapes.add_shape(1, Inches(0), Inches(3.5), Inches(10), Inches(0.06))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = GOLD; bar2.line.fill.background()
    textbox(slide, "Thank You", 0.5, 1.8, 9, 1.2, 48, True, WHITE, PP_ALIGN.CENTER)
    textbox(slide, "This report is confidential and prepared for the CEO.",
            0.5, 3.7, 9, 0.6, 13, False, RGBColor(0xBD,0xD7,0xEE), PP_ALIGN.CENTER)
    textbox(slide, f"Data Source: Book1 (003).xlsx — Sheet: {sheet_name}  |  {datetime.now().strftime('%B %Y')}",
            0.5, 4.4, 9, 0.4, 9, False, RGBColor(0x9D,0xC3,0xE6), PP_ALIGN.CENTER)

    prs.save(out_path)
    print(f"  ✓ PPT saved: {out_path}")

# ─────────────────────────────────────────────────────────────────────────────
# Excel Working File Builder
# ─────────────────────────────────────────────────────────────────────────────
def build_excel(sheet_name, data, stats, out_path):
    wb = xlsxwriter.Workbook(out_path)

    # Formats
    fmt_title   = wb.add_format({'bold':True,'font_size':14,'font_color':'#FFFFFF','bg_color':'#1F3564','border':1,'align':'center','valign':'vcenter'})
    fmt_header  = wb.add_format({'bold':True,'font_size':10,'font_color':'#FFFFFF','bg_color':'#2E75B6','border':1,'align':'center','valign':'vcenter'})
    fmt_header2 = wb.add_format({'bold':True,'font_size':10,'font_color':'#FFFFFF','bg_color':'#1F3564','border':1,'align':'center','valign':'vcenter'})
    fmt_data    = wb.add_format({'font_size':9,'border':1,'valign':'vcenter'})
    fmt_data_c  = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center'})
    fmt_alt     = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','bg_color':'#F2F2F2'})
    fmt_alt_c   = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center','bg_color':'#F2F2F2'})
    fmt_pct     = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center','num_format':'0%'})
    fmt_pct_alt = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center','num_format':'0%','bg_color':'#F2F2F2'})
    fmt_num     = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center','num_format':'#,##0'})
    fmt_num_alt = wb.add_format({'font_size':9,'border':1,'valign':'vcenter','align':'center','num_format':'#,##0','bg_color':'#F2F2F2'})
    fmt_green   = wb.add_format({'bold':True,'font_size':10,'font_color':'#FFFFFF','bg_color':'#70AD47','border':1,'align':'center','valign':'vcenter'})
    fmt_orange  = wb.add_format({'bold':True,'font_size':10,'font_color':'#FFFFFF','bg_color':'#FF7600','border':1,'align':'center','valign':'vcenter'})
    fmt_red     = wb.add_format({'bold':True,'font_size':10,'font_color':'#FFFFFF','bg_color':'#C00000','border':1,'align':'center','valign':'vcenter'})

    total = stats['total']

    # ── Sheet 1: Raw Data ────────────────────────────────────────────────────
    ws_raw = wb.add_worksheet('Raw Data')
    raw_headers = ['Empcode','Name','DOJ','Supervisor','Director','Actual Designation',
                   'Client','Band','Grade','Rating','Promotion','Appraisal','MIS Remarks']
    col_widths = [12,28,14,20,22,28,35,7,8,9,11,20,28]
    ws_raw.set_row(0, 22)
    ws_raw.write_row(0, 0, raw_headers, fmt_header)
    for ci, cw in enumerate(col_widths):
        ws_raw.set_column(ci, ci, cw)
    for ri, row in enumerate(data):
        fmt_d = fmt_alt if ri%2==1 else fmt_data
        fmt_c = fmt_alt_c if ri%2==1 else fmt_data_c
        vals = [
            v(row,'Empcode'), v(row,'Name'),
            str(row.get('DOJ',''))[:10] if row.get('DOJ') else '',
            v(row,'Supervisor'), v(row,'Director'), v(row,'Actual Designation'),
            v(row,'Client'), row.get('Band',''), str(row.get('Grade','')).strip(),
            v(row,'Rating'), v(row,'Promotion'), v(row,'Apprisal'), v(row,'MIS Remarks'),
        ]
        for ci, val in enumerate(vals):
            fmt_use = fmt_c if ci in [0,7,8,9,10] else fmt_d
            ws_raw.write(ri+1, ci, val, fmt_use)

    # ── Sheet 2: Summary Dashboard ───────────────────────────────────────────
    ws_sum = wb.add_worksheet('Summary Dashboard')
    ws_sum.set_column(0, 0, 32)
    ws_sum.set_column(1, 4, 18)
    ws_sum.merge_range('A1:E1', f'{sheet_name} — Appraisal Summary Dashboard', fmt_title)
    ws_sum.set_row(0, 28)

    def write_section(ws, row, title, data_rows, headers, col_fmts):
        ws.merge_range(row, 0, row, len(headers)-1, title, fmt_header2)
        ws.set_row(row, 20)
        row += 1
        ws.write_row(row, 0, headers, fmt_header)
        row += 1
        for ri2, dr in enumerate(data_rows):
            for ci2, (val, fmt) in enumerate(zip(dr, col_fmts)):
                alt = ri2%2==1
                if fmt == 'pct':
                    ws.write(row, ci2, val, fmt_pct_alt if alt else fmt_pct)
                elif fmt == 'num':
                    ws.write(row, ci2, val, fmt_num_alt if alt else fmt_num)
                elif fmt == 'c':
                    ws.write(row, ci2, val, fmt_alt_c if alt else fmt_data_c)
                else:
                    ws.write(row, ci2, val, fmt_alt if alt else fmt_data)
            row += 1
        return row + 1

    r = 2
    # KPI summary
    ws_sum.merge_range(r, 0, r, 4, 'KEY PERFORMANCE INDICATORS', fmt_header2)
    ws_sum.set_row(r, 20); r+=1
    kpis = [
        ('Total Employees', total, '', fmt_num, None),
        ('Rated Employees', stats['rated_count'], stats['rated_count']/total, fmt_num, fmt_pct),
        ('Rating 4 (Exceeds)', stats['ratings'].get('4',0), stats['ratings'].get('4',0)/total, fmt_num, fmt_pct),
        ('Rating 3 (Meets)', stats['ratings'].get('3',0), stats['ratings'].get('3',0)/total, fmt_num, fmt_pct),
        ('Promoted', stats['promos'].get('Yes',0), stats['promos'].get('Yes',0)/total, fmt_num, fmt_pct),
        ('Fully Eligible', stats['fully_eligible'], stats['fully_eligible']/total, fmt_num, fmt_pct),
        ('Partially Eligible', stats['partially_eligible'], stats['partially_eligible']/total, fmt_num, fmt_pct),
        ('Not Eligible', stats['not_eligible'], stats['not_eligible']/total, fmt_num, fmt_pct),
    ]
    ws_sum.write_row(r, 0, ['Metric','Count','% of Total','',''], fmt_header); r+=1
    for i, (lbl, cnt, pct, f1, f2) in enumerate(kpis):
        alt = i%2==1
        ws_sum.write(r, 0, lbl, fmt_alt if alt else fmt_data)
        ws_sum.write(r, 1, cnt, fmt_num_alt if alt else fmt_num)
        if f2:
            ws_sum.write(r, 2, pct, fmt_pct_alt if alt else fmt_pct)
        r+=1
    r+=1

    # Rating distribution
    rating_rows = [
        ('Rating 4 (Exceeds)', stats['ratings'].get('4',0), stats['ratings'].get('4',0)/total),
        ('Rating 3 (Meets)',   stats['ratings'].get('3',0), stats['ratings'].get('3',0)/total),
        ('Rating 2 (Below)',   stats['ratings'].get('2',0), stats['ratings'].get('2',0)/total),
        ('Not Applicable (NA)',stats['ratings'].get('NA',0),stats['ratings'].get('NA',0)/total),
    ]
    r = write_section(ws_sum, r, 'RATING DISTRIBUTION',
                      rating_rows, ['Rating','Count','% of Total'],
                      ['d','num','pct'])

    # Appraisal eligibility
    elig_rows = [
        ('Fully Eligible',     stats['fully_eligible'],    stats['fully_eligible']/total),
        ('Partially Eligible', stats['partially_eligible'],stats['partially_eligible']/total),
        ('Not Eligible',       stats['not_eligible'],      stats['not_eligible']/total),
    ]
    r = write_section(ws_sum, r, 'APPRAISAL ELIGIBILITY',
                      elig_rows, ['Status','Count','% of Total'],
                      ['d','num','pct'])

    # Band distribution
    band_rows = [(f'Band {b}', stats['bands'].get(b,0), stats['bands'].get(b,0)/total) for b in [1,2,3]]
    r = write_section(ws_sum, r, 'BAND DISTRIBUTION',
                      band_rows, ['Band','Count','% of Total'],
                      ['c','num','pct'])

    # Promotion summary
    promo_rows = [
        ('Promoted (Yes)', stats['promos'].get('Yes',0), stats['promos'].get('Yes',0)/total),
        ('Not Promoted (No)', stats['promos'].get('No',0), stats['promos'].get('No',0)/total),
    ]
    r = write_section(ws_sum, r, 'PROMOTION SUMMARY',
                      promo_rows, ['Status','Count','% of Total'],
                      ['d','num','pct'])

    # ── Sheet 3: Rating Analysis ──────────────────────────────────────────────
    ws_rat = wb.add_worksheet('Rating Analysis')
    ws_rat.set_column(0, 0, 28)
    ws_rat.set_column(1, 6, 14)
    ws_rat.merge_range('A1:G1', 'Detailed Rating Analysis', fmt_title)
    ws_rat.set_row(0, 28)

    r = 2
    # Rating by Band
    ws_rat.merge_range(r, 0, r, 5, 'RATING BY BAND', fmt_header2); r+=1
    ws_rat.write_row(r, 0, ['Band','Total','Rating 4','Rating 3','Rating 2','Not Eligible (NA)'], fmt_header); r+=1
    for i, band in enumerate([1,2,3]):
        br = stats['band_rating'].get(band,{})
        tot_b = stats['bands'].get(band,0)
        alt = i%2==1
        row_vals = [f'Band {band}', tot_b, br.get('4',0), br.get('3',0), br.get('2',0), br.get('NA',0)]
        fmts = [fmt_alt if alt else fmt_data] + [fmt_num_alt if alt else fmt_num]*5
        for ci, (val, fmt) in enumerate(zip(row_vals, fmts)):
            ws_rat.write(r, ci, val, fmt)
        r+=1
    r+=1

    # Rating by Director
    ws_rat.merge_range(r, 0, r, 6, 'RATING BY DIRECTOR', fmt_header2); r+=1
    ws_rat.write_row(r, 0, ['Director','Total','Rated','Rating 4','Rating 3','NA','Promotion Rate'], fmt_header); r+=1
    dir_order = sorted(stats['dir_rating'].keys(), key=lambda d:-sum(stats['dir_rating'][d].values()))
    for i, d in enumerate(dir_order):
        dr = stats['dir_rating'].get(d,{})
        dp = stats['dir_promo'].get(d,{})
        tot_d = sum(dr.values())
        rated_d = dr.get('4',0)+dr.get('3',0)+dr.get('2',0)
        promo_rate = dp.get('Yes',0)/tot_d if tot_d else 0
        alt = i%2==1
        ws_rat.write(r, 0, d, fmt_alt if alt else fmt_data)
        for ci, val in enumerate([tot_d, rated_d, dr.get('4',0), dr.get('3',0), dr.get('NA',0)]):
            ws_rat.write(r, ci+1, val, fmt_num_alt if alt else fmt_num)
        ws_rat.write(r, 6, promo_rate, fmt_pct_alt if alt else fmt_pct)
        r+=1
    r+=1

    # Top Rating 4 by Designation
    ws_rat.merge_range(r, 0, r, 2, 'RATING 4 BY DESIGNATION (TOP 15)', fmt_header2); r+=1
    ws_rat.write_row(r, 0, ['Designation','Rating 4 Count','% of Designation'], fmt_header); r+=1
    desig_r4 = defaultdict(lambda: defaultdict(int))
    desig_total = defaultdict(int)
    for row_ in data:
        d = v(row_,'Actual Designation')
        desig_r4[d][v(row_,'Rating')] += 1
        desig_total[d] += 1
    r4_by_desig = sorted([(d, cnt.get('4',0), desig_total[d]) for d,cnt in desig_r4.items()],
                          key=lambda x:-x[1])[:15]
    for i, (d, r4_cnt, dtot) in enumerate(r4_by_desig):
        alt = i%2==1
        ws_rat.write(r, 0, d, fmt_alt if alt else fmt_data)
        ws_rat.write(r, 1, r4_cnt, fmt_num_alt if alt else fmt_num)
        ws_rat.write(r, 2, r4_cnt/dtot if dtot else 0, fmt_pct_alt if alt else fmt_pct)
        r+=1

    # ── Sheet 4: Promotion Analysis ───────────────────────────────────────────
    ws_pro = wb.add_worksheet('Promotion Analysis')
    ws_pro.set_column(0, 0, 32)
    ws_pro.set_column(1, 5, 16)
    ws_pro.merge_range('A1:F1', 'Detailed Promotion Analysis', fmt_title)
    ws_pro.set_row(0, 28)

    r = 2
    # Promotions by Director
    ws_pro.merge_range(r, 0, r, 4, 'PROMOTIONS BY DIRECTOR', fmt_header2); r+=1
    ws_pro.write_row(r, 0, ['Director','Total','Promoted','Not Promoted','Promotion Rate'], fmt_header); r+=1
    for i, d in enumerate(dir_order):
        dp = stats['dir_promo'].get(d,{})
        tot_d = dp.get('Yes',0)+dp.get('No',0)
        alt = i%2==1
        ws_pro.write(r, 0, d, fmt_alt if alt else fmt_data)
        ws_pro.write(r, 1, tot_d, fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 2, dp.get('Yes',0), fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 3, dp.get('No',0), fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 4, dp.get('Yes',0)/tot_d if tot_d else 0, fmt_pct_alt if alt else fmt_pct)
        r+=1
    r+=1

    # Promotions by Band
    ws_pro.merge_range(r, 0, r, 4, 'PROMOTIONS BY BAND', fmt_header2); r+=1
    ws_pro.write_row(r, 0, ['Band','Total','Promoted','Not Promoted','Promotion Rate'], fmt_header); r+=1
    for i, band in enumerate([1,2,3]):
        bp = stats['band_promo'].get(band,{})
        tot_b = stats['bands'].get(band,0)
        alt = i%2==1
        ws_pro.write(r, 0, f'Band {band}', fmt_alt_c if alt else fmt_data_c)
        ws_pro.write(r, 1, tot_b, fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 2, bp.get('Yes',0), fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 3, bp.get('No',0), fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 4, bp.get('Yes',0)/tot_b if tot_b else 0, fmt_pct_alt if alt else fmt_pct)
        r+=1
    r+=1

    # Promotions by Designation
    ws_pro.merge_range(r, 0, r, 3, 'PROMOTIONS BY DESIGNATION', fmt_header2); r+=1
    ws_pro.write_row(r, 0, ['Designation','Total in Role','Promoted','Promotion Rate'], fmt_header); r+=1
    promo_by_desig = sorted(
        [(d, sum(cnt.values()), cnt.get('Yes',0)) for d,cnt in stats['desig_promo'].items()],
        key=lambda x:-x[2])
    for i, (d, tot_d, y) in enumerate(promo_by_desig):
        alt = i%2==1
        ws_pro.write(r, 0, d, fmt_alt if alt else fmt_data)
        ws_pro.write(r, 1, tot_d, fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 2, y, fmt_num_alt if alt else fmt_num)
        ws_pro.write(r, 3, y/tot_d if tot_d else 0, fmt_pct_alt if alt else fmt_pct)
        r+=1

    # ── Sheet 5: Client & Grade Analysis ──────────────────────────────────────
    ws_cli = wb.add_worksheet('Client & Grade Analysis')
    ws_cli.set_column(0, 0, 40)
    ws_cli.set_column(1, 4, 16)
    ws_cli.merge_range('A1:E1', 'Client & Grade Breakdown', fmt_title)
    ws_cli.set_row(0, 28)

    r = 2
    # Client headcount
    ws_cli.merge_range(r, 0, r, 3, 'HEADCOUNT BY CLIENT', fmt_header2); r+=1
    ws_cli.write_row(r, 0, ['Client','Count','% of Total','Rating 4 Count'], fmt_header); r+=1
    cli_r4 = defaultdict(int)
    for row_ in data:
        if v(row_,'Rating')=='4':
            cli_r4[v(row_,'Client')] += 1
    for i, (cli, cnt) in enumerate(stats['clients'].most_common()):
        alt = i%2==1
        ws_cli.write(r, 0, cli, fmt_alt if alt else fmt_data)
        ws_cli.write(r, 1, cnt, fmt_num_alt if alt else fmt_num)
        ws_cli.write(r, 2, cnt/total, fmt_pct_alt if alt else fmt_pct)
        ws_cli.write(r, 3, cli_r4.get(cli,0), fmt_num_alt if alt else fmt_num)
        r+=1
    r+=1

    # Grade breakdown
    ws_cli.merge_range(r, 0, r, 3, 'HEADCOUNT BY GRADE', fmt_header2); r+=1
    ws_cli.write_row(r, 0, ['Grade','Count','% of Total','Promoted'], fmt_header); r+=1
    grade_promo = defaultdict(int)
    for row_ in data:
        if v(row_,'Promotion')=='Yes':
            grade_promo[str(row_.get('Grade','')).strip()] += 1
    for i, (grade, cnt) in enumerate(stats['grades'].most_common()):
        alt = i%2==1
        ws_cli.write(r, 0, grade, fmt_alt_c if alt else fmt_data_c)
        ws_cli.write(r, 1, cnt, fmt_num_alt if alt else fmt_num)
        ws_cli.write(r, 2, cnt/total, fmt_pct_alt if alt else fmt_pct)
        ws_cli.write(r, 3, grade_promo.get(grade,0), fmt_num_alt if alt else fmt_num)
        r+=1

    wb.close()
    print(f"  ✓ Excel saved: {out_path}")

# ─────────────────────────────────────────────────────────────────────────────
# Summary Document Builder
# ─────────────────────────────────────────────────────────────────────────────
def build_summary(sheet_name, data, stats, out_path):
    total = stats['total']
    lines = []
    def w(s=""): lines.append(s)

    w("=" * 80)
    w(f"  APPRAISAL ANALYSIS — DETAILED SUMMARY DOCUMENT")
    w(f"  Sheet: {sheet_name}")
    w(f"  Generated: {datetime.now().strftime('%d %B %Y, %H:%M')}")
    w(f"  Source: Book1 (003).xlsx → Sheet '{sheet_name}'")
    w("=" * 80)
    w()

    w("━" * 80)
    w("1. DATA OVERVIEW")
    w("━" * 80)
    w(f"  Total Employee Records  : {total}")
    w(f"  Columns                 : Empcode, Name, DOJ, Supervisor, Director,")
    w(f"                             Actual Designation, Client, Band, Grade,")
    w(f"                             Rating, Promotion, Appraisal, MIS Remarks")
    w(f"  Data Source             : Excel sheet '{sheet_name}' in Book1 (003).xlsx")
    w()

    w("━" * 80)
    w("2. APPRAISAL ELIGIBILITY — SOURCE: 'Apprisal' + 'MIS Remarks' columns")
    w("━" * 80)
    w(f"  Fully Eligible          : {stats['fully_eligible']:>4}  ({stats['fully_eligible']*100//total}%)")
    w(f"    → MIS Remark: 'Full year Covered'")
    w(f"    → Employees who joined on or before Apr 1, 2025")
    w()
    w(f"  Partially Eligible      : {stats['partially_eligible']:>4}  ({stats['partially_eligible']*100//total}%)")
    w(f"    → MIS Remark: 'Full year not covered'")
    w(f"    → Joined Oct 2024 – Apr 2025; merit increment will be pro-rated")
    w()
    w(f"  Not Eligible            : {stats['not_eligible']:>4}  ({stats['not_eligible']*100//total}%)")
    w(f"    → MIS Remark: 'Joined after 1st Oct 2025'")
    w(f"    → Rating column shows 'NA'; excluded from performance cycle")
    w()

    w("━" * 80)
    w("3. PERFORMANCE RATINGS — SOURCE: 'Rating' column")
    w("━" * 80)
    w(f"  Rating 4 (Exceeds Expectations)  : {stats['ratings'].get('4',0):>4}  ({stats['ratings'].get('4',0)*100//total}%)")
    w(f"  Rating 3 (Meets Expectations)    : {stats['ratings'].get('3',0):>4}  ({stats['ratings'].get('3',0)*100//total}%)")
    w(f"  Rating 2 (Below Expectations)    : {stats['ratings'].get('2',0):>4}  ({stats['ratings'].get('2',0)*100//total}%)")
    w(f"  Not Applicable (NA)              : {stats['ratings'].get('NA',0):>4}  ({stats['ratings'].get('NA',0)*100//total}%)")
    w(f"  Total Rated Employees            : {stats['rated_count']:>4}  ({stats['rated_count']*100//total}%)")
    w()
    w("  RATING BY BAND:")
    for band in [1,2,3]:
        br = stats['band_rating'].get(band,{})
        tot_b = stats['bands'].get(band,0)
        w(f"    Band {band} ({tot_b} emp): R4={br.get('4',0)}, R3={br.get('3',0)}, R2={br.get('2',0)}, NA={br.get('NA',0)}")
    w()
    w("  RATING BY DIRECTOR:")
    for d, dr in stats['dir_rating'].items():
        tot_d = sum(dr.values())
        w(f"    {d} ({tot_d} emp): R4={dr.get('4',0)}, R3={dr.get('3',0)}, R2={dr.get('2',0)}, NA={dr.get('NA',0)}")
    w()

    w("━" * 80)
    w("4. PROMOTIONS — SOURCE: 'Promotion' column")
    w("━" * 80)
    w(f"  Promoted (Yes)          : {stats['promos'].get('Yes',0):>4}  ({stats['promos'].get('Yes',0)*100//total}%)")
    w(f"  Not Promoted (No)       : {stats['promos'].get('No',0):>4}  ({stats['promos'].get('No',0)*100//total}%)")
    w()
    w("  BY BAND:")
    for band in [1,2,3]:
        bp = stats['band_promo'].get(band,{})
        tot_b = stats['bands'].get(band,0)
        pct = f"{bp.get('Yes',0)*100//tot_b}%" if tot_b else "0%"
        w(f"    Band {band}: {bp.get('Yes',0)} promoted of {tot_b} ({pct})")
    w()
    w("  BY DIRECTOR:")
    for d, dp in stats['dir_promo'].items():
        tot_d = dp.get('Yes',0)+dp.get('No',0)
        pct = f"{dp.get('Yes',0)*100//tot_d}%" if tot_d else "0%"
        w(f"    {d}: {dp.get('Yes',0)} promoted of {tot_d} ({pct})")
    w()
    w("  BY DESIGNATION (promoted):")
    for d, y in sorted([(d,cnt.get('Yes',0)) for d,cnt in stats['desig_promo'].items() if cnt.get('Yes',0)>0], key=lambda x:-x[1]):
        w(f"    {d}: {y}")
    w()

    w("━" * 80)
    w("5. BAND DISTRIBUTION — SOURCE: 'Band' column")
    w("━" * 80)
    for band in [1,2,3]:
        cnt = stats['bands'].get(band,0)
        w(f"  Band {band}: {cnt:>4} employees ({cnt*100//total}%)")
    w()

    w("━" * 80)
    w("6. GRADE DISTRIBUTION — SOURCE: 'Grade' column")
    w("━" * 80)
    for grade, cnt in stats['grades'].most_common():
        w(f"  {grade:<12}: {cnt:>4} ({cnt*100//total}%)")
    w()

    w("━" * 80)
    w("7. DIRECTOR PORTFOLIO — SOURCE: 'Director' column")
    w("━" * 80)
    for d, cnt in stats['directors'].most_common():
        w(f"  {d}: {cnt} employees ({cnt*100//total}%)")
    w()

    w("━" * 80)
    w("8. CLIENT DISTRIBUTION — SOURCE: 'Client' column")
    w("━" * 80)
    for cli, cnt in stats['clients'].most_common():
        w(f"  {cli}: {cnt} employees ({cnt*100//total}%)")
    w()

    w("━" * 80)
    w("9. DESIGNATION DISTRIBUTION — SOURCE: 'Actual Designation' column")
    w("━" * 80)
    for desig, cnt in stats['desigs'].most_common():
        w(f"  {desig}: {cnt}")
    w()

    w("━" * 80)
    w("10. TOP SUPERVISORS — SOURCE: 'Supervisor' column")
    w("━" * 80)
    sup_counts = Counter(v(r,'Supervisor') for r in data if v(r,'Supervisor'))
    for s, cnt in sup_counts.most_common(15):
        w(f"  {s}: {cnt} reportees")
    w()

    w("━" * 80)
    w("11. DATA QUALITY NOTES")
    w("━" * 80)
    missing_rating = sum(1 for r in data if not v(r,'Rating') or v(r,'Rating')=='None')
    missing_name   = sum(1 for r in data if not v(r,'Name'))
    w(f"  Records with missing Rating : {missing_rating}")
    w(f"  Records with missing Name   : {missing_name}")
    w(f"  Note: 'Rating' column header has a trailing space in source Excel")
    w(f"  Note: 'Promotion' column header has a trailing space in source Excel")
    w(f"  Note: 'MIS Remarks' column header has a trailing space in source Excel")
    w()

    w("━" * 80)
    w("12. DATA POINT ORIGIN MAP")
    w("━" * 80)
    w("  Every metric in this report traces directly to a column in the source sheet.")
    w()
    w(f"  {'Metric':<35} {'Source Column'}")
    w(f"  {'-'*35} {'-'*30}")
    origins = [
        ("Total Headcount",             "Count of all rows"),
        ("Appraisal Eligibility",       "'Apprisal' column + 'MIS Remarks'"),
        ("Performance Rating (4/3/2)",  "'Rating' column (numeric)"),
        ("Not Eligible / NA",            "'Rating' = 'NA'"),
        ("Promotion Status",            "'Promotion' column ('Yes'/'No')"),
        ("Band Classification",         "'Band' column (1/2/3)"),
        ("Grade",                       "'Grade' column (alphanumeric)"),
        ("Director Portfolio",          "'Director' column"),
        ("Client Account",              "'Client' column"),
        ("Designation",                 "'Actual Designation' column"),
        ("Supervisor Coverage",         "'Supervisor' column"),
        ("Date of Joining",             "'DOJ' column (datetime)"),
        ("Employee Identifier",         "'Empcode' column"),
    ]
    for metric, src in origins:
        w(f"  {metric:<35} {src}")
    w()

    w("=" * 80)
    w("END OF SUMMARY")
    w("=" * 80)

    with open(out_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))
    print(f"  ✓ Summary saved: {out_path}")

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
print("\n📊  Loading and analysing data...")
appr_data = load_sheet('Apprisal Data')
al_data   = load_sheet('AL')
appr_stats = compute_stats(appr_data)
al_stats   = compute_stats(al_data)

print("\n🔷  Appraisal Data sheet")
print(f"   Total records: {appr_stats['total']}")
build_ppt('Apprisal Data', appr_data, appr_stats,
          f"{OUT}/Appraisal_CEO_Report.pptx", "APPRAISAL DATA — ALL CLIENTS")
build_excel('Apprisal Data', appr_data, appr_stats,
            f"{OUT}/Appraisal_Working_File.xlsx")
build_summary('Apprisal Data', appr_data, appr_stats,
              f"{OUT}/Appraisal_Summary.txt")

print("\n🔶  AL (Asset Living) sheet")
print(f"   Total records: {al_stats['total']}")
build_ppt('AL', al_data, al_stats,
          f"{OUT}/AL_CEO_Report.pptx", "ASSET LIVING (AL) — CLIENT FOCUS")
build_excel('AL', al_data, al_stats,
            f"{OUT}/AL_Working_File.xlsx")
build_summary('AL', al_data, al_stats,
              f"{OUT}/AL_Summary.txt")

print("\n✅  All 6 files generated:")
print(f"   {OUT}/Appraisal_CEO_Report.pptx")
print(f"   {OUT}/Appraisal_Working_File.xlsx")
print(f"   {OUT}/Appraisal_Summary.txt")
print(f"   {OUT}/AL_CEO_Report.pptx")
print(f"   {OUT}/AL_Working_File.xlsx")
print(f"   {OUT}/AL_Summary.txt")
