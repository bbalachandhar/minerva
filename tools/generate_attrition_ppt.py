from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────
DARK_NAVY   = RGBColor(0x0D, 0x1B, 0x3E)   # slide BG / headers
MID_BLUE    = RGBColor(0x1A, 0x53, 0x9E)   # accent bar
ACCENT_BLUE = RGBColor(0x27, 0x7D, 0xC2)   # lighter accent
RED_ALERT   = RGBColor(0xC0, 0x39, 0x2B)   # warning numbers
GREEN_OK    = RGBColor(0x1E, 0x8B, 0x4C)   # positive numbers
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GOLD        = RGBColor(0xF3, 0x9C, 0x12)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ══════════════════════════════════════════════════════════════════
#  HELPERS
# ══════════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_width
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width if line_width else Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, l, t, w, h, text, font_size=14, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_slide_bg(slide, color=DARK_NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.35, fill=MID_BLUE)
    add_textbox(slide, 0.35, 0.12, 12.5, 0.65, title,
                font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 0.35, 0.75, 12.5, 0.45, subtitle,
                    font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF))


def add_footer(slide, text="Confidential  |  HR / People Analytics  |  31 March 2026"):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=RGBColor(0x0A, 0x14, 0x2A))
    add_textbox(slide, 0.3, 7.12, 12.5, 0.3, text, font_size=9,
                color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.CENTER)


def kpi_box(slide, l, t, w, h, value, label, val_color=WHITE, bg=MID_BLUE):
    add_rect(slide, l, t, w, h, fill=bg, line_color=ACCENT_BLUE, line_width=Pt(1.5))
    add_textbox(slide, l, t + 0.1, w, 0.65, value,
                font_size=32, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, l, t + 0.7, w, 0.4, label,
                font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.CENTER)


def bar_chart(slide, l, t, w, h, data, max_val, bar_color=ACCENT_BLUE, label_color=WHITE):
    """data = list of (label, value) tuples"""
    n = len(data)
    row_h = h / n
    for i, (lbl, val) in enumerate(data):
        y = t + i * row_h
        # label
        add_textbox(slide, l, y + row_h * 0.1, 2.0, row_h * 0.8,
                    lbl, font_size=10, color=label_color)
        # bar bg
        bar_x = l + 2.1
        bar_w = w - 2.1 - 0.7
        add_rect(slide, bar_x, y + row_h * 0.15,
                 bar_w, row_h * 0.65, fill=RGBColor(0x22, 0x33, 0x55))
        filled = bar_w * (val / max_val)
        add_rect(slide, bar_x, y + row_h * 0.15,
                 filled, row_h * 0.65, fill=bar_color)
        # value label
        add_textbox(slide, bar_x + filled + 0.05, y + row_h * 0.1,
                    0.55, row_h * 0.8, str(val),
                    font_size=10, bold=True, color=label_color)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)

# decorative left strip
add_rect(slide, 0, 0, 0.6, 7.5, fill=ACCENT_BLUE)
add_rect(slide, 0.6, 0, 0.12, 7.5, fill=GOLD)

# main title block
add_rect(slide, 1.0, 1.6, 11.0, 1.0, fill=MID_BLUE)
add_textbox(slide, 1.2, 1.65, 10.5, 0.9,
            "ATTRITION DATA ANALYSIS REPORT",
            font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 2.85, 11.0, 0.55,
            "Asset Living, LLC  |  F&A – Operations  |  July 2025 – March 2026",
            font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(slide, 1.0, 3.45, 11.0, 0.45,
            "Generated: 31 March 2026  |  HR / People Analytics",
            font_size=13, color=RGBColor(0x99, 0xAA, 0xCC), align=PP_ALIGN.CENTER,
            italic=True)

# KPI strip
for i, (val, lbl, col) in enumerate([
    ("108", "Total Exits", RED_ALERT),
    ("9 Months", "Period Covered", WHITE),
    ("17 / mo", "2026 Run Rate", RED_ALERT),
    ("85 %", "Undesired Exits", RED_ALERT),
]):
    x = 1.2 + i * 2.75
    kpi_box(slide, x, 4.35, 2.4, 1.25, val, lbl, val_color=col)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 2 — EXECUTIVE SNAPSHOT
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "Executive Snapshot", "108 Exits  |  Jul 2025 – Mar 2026  |  Asset Living F&A")

# 4 big KPI cards
kpis = [
    ("108",    "Total Attritions",     RED_ALERT,  RGBColor(0x12, 0x2A, 0x5E)),
    ("17/mo",  "2026 Monthly Rate",    RED_ALERT,  RGBColor(0x12, 0x2A, 0x5E)),
    ("85%",    "Undesired Exits",      GOLD,       RGBColor(0x12, 0x2A, 0x5E)),
    ("+79%",   "Rate vs H2 2025",      RED_ALERT,  RGBColor(0x12, 0x2A, 0x5E)),
]
for i, (val, lbl, vc, bg) in enumerate(kpis):
    x = 0.35 + i * 3.22
    kpi_box(slide, x, 1.6, 2.9, 1.3, val, lbl, val_color=vc, bg=bg)

# Attrition type breakdown
add_rect(slide, 0.35, 3.1, 5.8, 3.55, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 0.55, 3.18, 5.4, 0.4, "ATTRITION TYPE BREAKDOWN",
            font_size=12, bold=True, color=ACCENT_BLUE)
bar_chart(slide, 0.45, 3.65, 5.6, 2.8,
          [("Resignation", 82), ("Abscond", 21), ("Termination", 5)], 90)

# Top reasons
add_rect(slide, 6.45, 3.1, 6.55, 3.55, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 6.65, 3.18, 6.1, 0.4, "TOP ATTRITION REASONS",
            font_size=12, bold=True, color=ACCENT_BLUE)
bar_chart(slide, 6.55, 3.65, 6.3, 2.8,
          [("Personal Reason", 39), ("Health Reason", 31),
           ("Other Opportunity", 14), ("Disciplinary", 8),
           ("Higher Studies", 7)], 45)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 3 — 2025 vs 2026 YEAR COMPARISON
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "Year-Wise Comparison: 2025 vs 2026",
               "H2 2025 (Jul–Dec)  vs  Q1 2026 (Jan–Mar)")

# Comparison table
headers  = ["Metric", "2025 (Jul–Dec)", "2026 (Jan–Mar)", "Trend"]
rows = [
    ["Total Attritions",   "57",      "51",        "—"],
    ["Months Covered",     "6",       "3",         "—"],
    ["Monthly Average",    "9.5/mo",  "17.0/mo",   "▲ +79%"],
    ["Resignations",       "40",      "42",        "▲"],
    ["Absconds",           "14",      "7",         "▼ Improved"],
    ["Terminations",       "3",       "5",         "▲"],
    ["Undesired Exits",    "52 (91%)", "40 (78%)", "▼ Improved %"],
    ["Desired Exits",      "5 (9%)",  "11 (22%)",  "▲ Improved"],
]

col_x = [0.3, 3.7, 7.1, 10.5]
col_w = [3.3, 3.3, 3.3, 2.6]
row_h  = 0.38
table_t = 1.55

# header row
add_rect(slide, 0.3, table_t, 12.7, row_h, fill=ACCENT_BLUE)
for j, hdr in enumerate(headers):
    add_textbox(slide, col_x[j] + 0.05, table_t + 0.04, col_w[j], row_h - 0.08,
                hdr, font_size=11, bold=True, color=WHITE)

for i, row in enumerate(rows):
    y = table_t + (i + 1) * row_h
    bg = RGBColor(0x10, 0x22, 0x48) if i % 2 == 0 else RGBColor(0x0D, 0x1B, 0x3E)
    add_rect(slide, 0.3, y, 12.7, row_h, fill=bg)
    for j, cell in enumerate(row):
        col = WHITE
        if j == 3:
            if "▲" in cell and "Improved" not in cell:
                col = RED_ALERT
            elif "▼" in cell or "Improved" in cell:
                col = GREEN_OK
        add_textbox(slide, col_x[j] + 0.05, y + 0.04, col_w[j], row_h - 0.08,
                    cell, font_size=10, color=col)

# Monthly trend bars
add_rect(slide, 0.3, 5.1, 12.7, 2.0, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 0.5, 5.15, 12.0, 0.35, "MONTHLY TREND",
            font_size=12, bold=True, color=ACCENT_BLUE)

months = [("Jul", 6), ("Aug", 9), ("Sep", 8), ("Oct", 8), ("Nov", 8),
          ("Dec", 18), ("Jan", 13), ("Feb", 22), ("Mar", 16)]
bar_w  = 1.08
bar_gap = 0.25
max_v  = 22
chart_l = 0.55
chart_t = 5.55
chart_h = 1.3

for idx, (mon, val) in enumerate(months):
    x = chart_l + idx * (bar_w + bar_gap)
    bh = chart_h * (val / max_v)
    by = chart_t + chart_h - bh
    col = RED_ALERT if val >= 18 else (GOLD if val >= 13 else ACCENT_BLUE)
    # black bg column
    add_rect(slide, x, chart_t, bar_w, chart_h, fill=RGBColor(0x07, 0x10, 0x22))
    add_rect(slide, x, by, bar_w, bh, fill=col)
    add_textbox(slide, x, by - 0.28, bar_w, 0.28, str(val),
                font_size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
    lbl_col = RGBColor(0xAA, 0xBB, 0xCC) if idx < 6 else WHITE
    add_textbox(slide, x, chart_t + chart_h + 0.03, bar_w, 0.22, mon,
                font_size=9, color=lbl_col, align=PP_ALIGN.CENTER)

# year labels
add_textbox(slide, chart_l, chart_t + chart_h + 0.27, 6.85 * 1.0, 0.22,
            "◀  2025 (H2)", font_size=9, color=RGBColor(0x99, 0xAA, 0xCC))
add_textbox(slide, chart_l + 7.2, chart_t + chart_h + 0.27, 4.0, 0.22,
            "2026 (Q1)  ▶", font_size=9, color=WHITE)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 4 — TENURE ANALYSIS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "Tenure Analysis",
               "When are people leaving? — 57% exit within first 12 months")

tenure_data = [
    ("0 – 30 days",    2,   "1.9%"),
    ("31 – 60 days",   9,   "8.3%"),
    ("61 – 90 days",   7,   "6.5%"),
    ("91 – 180 days",  34,  "31.5%  ← HIGHEST RISK"),
    ("181 – 365 days", 16,  "14.8%"),
    ("1 – 2 Years",    28,  "25.9%"),
    ("2 – 3 Years",    12,  "11.1%"),
]

max_t = 34
row_h = 0.57
t_start = 1.55

for i, (lbl, val, pct) in enumerate(tenure_data):
    y = t_start + i * row_h
    is_peak = val == 34
    bg = RGBColor(0x4A, 0x10, 0x10) if is_peak else (
         RGBColor(0x10, 0x22, 0x48) if i % 2 == 0 else RGBColor(0x0D, 0x1B, 0x3E))
    add_rect(slide, 0.3, y, 12.7, row_h - 0.04, fill=bg)

    # label
    add_textbox(slide, 0.45, y + 0.08, 2.3, row_h - 0.15,
                lbl, font_size=11, bold=is_peak,
                color=GOLD if is_peak else WHITE)
    # bar
    bar_area_l = 2.85
    bar_area_w = 7.5
    filled = bar_area_w * (val / max_t)
    add_rect(slide, bar_area_l, y + 0.1, bar_area_w, row_h - 0.25,
             fill=RGBColor(0x1A, 0x2A, 0x4A))
    add_rect(slide, bar_area_l, y + 0.1, filled, row_h - 0.25,
             fill=RED_ALERT if is_peak else ACCENT_BLUE)

    # count
    add_textbox(slide, bar_area_l + filled + 0.1, y + 0.08, 0.5, row_h - 0.15,
                str(val), font_size=11, bold=True,
                color=RED_ALERT if is_peak else WHITE)
    # pct
    add_textbox(slide, 11.0, y + 0.08, 2.0, row_h - 0.15,
                pct, font_size=10,
                color=GOLD if is_peak else RGBColor(0xCC, 0xDD, 0xFF))

# Insight box
add_rect(slide, 0.3, 5.6, 12.7, 1.5, fill=RGBColor(0x10, 0x28, 0x50))
add_rect(slide, 0.3, 5.6, 0.18, 1.5, fill=GOLD)
add_textbox(slide, 0.65, 5.68, 12.1, 0.35, "KEY INSIGHT",
            font_size=11, bold=True, color=GOLD)
add_textbox(slide, 0.65, 6.02, 12.1, 0.9,
            "57% of all attritions occur within the first 12 months. The 91–180 day window (3–6 months) "
            "is the single largest risk bucket with 34 exits (31% of total). Structured check-ins and "
            "buddy programs at this stage are critical to retention.",
            font_size=11, color=WHITE, wrap=True)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 5 — BAND & REASONS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "Band Distribution & Exit Reasons",
               "Who is leaving and why?")

# LEFT — Band distribution
add_rect(slide, 0.3, 1.55, 5.9, 5.6, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 0.5, 1.65, 5.5, 0.38, "BAND DISTRIBUTION",
            font_size=13, bold=True, color=ACCENT_BLUE)

bands = [("Band 1  (Entry/Junior)", 88, "81%", RED_ALERT),
         ("Band 2  (Mid Level)",     18, "17%", GOLD),
         ("Band 3  (Senior)",         2,  "2%", GREEN_OK)]
for i, (lbl, val, pct, col) in enumerate(bands):
    y = 2.2 + i * 1.5
    add_rect(slide, 0.5, y, 5.5, 1.2, fill=RGBColor(0x07, 0x14, 0x30))
    add_textbox(slide, 0.65, y + 0.08, 3.5, 0.4,
                lbl, font_size=12, bold=(i == 0), color=WHITE)
    bw = 4.5 * (val / 88)
    add_rect(slide, 0.65, y + 0.55, 4.5, 0.38, fill=RGBColor(0x1A, 0x2A, 0x4A))
    add_rect(slide, 0.65, y + 0.55, bw,  0.38, fill=col)
    add_textbox(slide, 5.1, y + 0.5, 0.8, 0.45,
                pct, font_size=16, bold=True, color=col, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.5, 6.55, 5.5, 0.5,
            "81% of exits are Band 1. Retention at entry level is the highest-impact lever.",
            font_size=10, color=RGBColor(0xCC, 0xDD, 0xFF), italic=True, wrap=True)

# RIGHT — Reasons
add_rect(slide, 6.55, 1.55, 6.45, 5.6, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 6.75, 1.65, 6.0, 0.38, "EXIT REASONS (Full Breakdown)",
            font_size=13, bold=True, color=ACCENT_BLUE)

reasons = [
    ("Personal Reason",   39, "36.1%"),
    ("Health Reason",     31, "28.7%"),
    ("Other Opportunity", 14, "13.0%"),
    ("Disciplinary",       8,  "7.4%"),
    ("Higher Studies",     7,  "6.5%"),
    ("Performance",        5,  "4.6%"),
    ("Work Related",       2,  "1.9%"),
    ("NCNS",               1,  "0.9%"),
    ("Relocation",         1,  "0.9%"),
]
r_row_h = 0.5
for i, (lbl, val, pct) in enumerate(reasons):
    y = 2.15 + i * r_row_h
    is_high = val >= 30
    add_textbox(slide, 6.75, y, 2.8, r_row_h - 0.06,
                lbl, font_size=10, color=WHITE)
    bw = 2.8 * (val / 39)
    add_rect(slide, 9.65, y + 0.07, 2.8, r_row_h - 0.2,
             fill=RGBColor(0x1A, 0x2A, 0x4A))
    add_rect(slide, 9.65, y + 0.07, bw,  r_row_h - 0.2,
             fill=RED_ALERT if is_high else ACCENT_BLUE)
    add_textbox(slide, 12.5, y, 0.4, r_row_h - 0.06,
                pct, font_size=9, color=RGBColor(0xCC, 0xDD, 0xFF))

add_textbox(slide, 6.75, 6.55, 6.0, 0.5,
            "Personal + Health = 65% of exits. Likely signals workload / WLB concerns.",
            font_size=10, color=RGBColor(0xCC, 0xDD, 0xFF), italic=True, wrap=True)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 6 — DESIRED vs UNDESIRED
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "Desired vs Undesired Exits",
               "Are we in control of who leaves?")

# Big numbers
for i, (val, lbl, sub, col, bg) in enumerate([
    ("92", "UNDESIRED", "Talent we wanted to keep — 85%", RED_ALERT, RGBColor(0x3A, 0x08, 0x08)),
    ("16", "DESIRED",   "Exits we initiated/accepted — 15%", GREEN_OK, RGBColor(0x08, 0x2A, 0x18)),
]):
    x = 0.5 + i * 6.55
    add_rect(slide, x, 1.6, 6.1, 3.2, fill=bg,
             line_color=col, line_width=Pt(2))
    add_textbox(slide, x, 1.75, 6.1, 1.3, val,
                font_size=72, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.0, 6.1, 0.55, lbl,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.5, 6.1, 0.45, sub,
                font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.CENTER, italic=True)

# Year comparison note
add_rect(slide, 0.3, 5.05, 12.7, 2.05, fill=RGBColor(0x10, 0x22, 0x48))
add_textbox(slide, 0.5, 5.13, 12.2, 0.38, "YEAR-ON-YEAR IMPROVEMENT",
            font_size=13, bold=True, color=ACCENT_BLUE)

items = [
    ("Undesired exits %",   "91% in 2025  →  78% in 2026",  "▼ Improving",  GREEN_OK),
    ("Desired exits %",     " 9% in 2025  →  22% in 2026",  "▲ Improving",  GREEN_OK),
    ("Absconds",            "14 in 2025   →   7 in 2026",   "▼ 50% reduction", GREEN_OK),
]
for i, (k, v, note, col) in enumerate(items):
    y = 5.55 + i * 0.47
    add_textbox(slide, 0.55, y, 3.5, 0.4, k, font_size=11, bold=True, color=WHITE)
    add_textbox(slide, 4.0, y, 5.5, 0.4, v, font_size=11, color=RGBColor(0xCC, 0xDD, 0xFF))
    add_textbox(slide, 9.6, y, 3.3, 0.4, note, font_size=11, bold=True, color=col)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 7 — 3 BIGGEST PROBLEMS
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "The 3 Biggest Problems",
               "Leadership must act on these signals now")

problems = [
    ("1",
     "EARLY EXITS DOMINATING",
     "57% of all exits happen within the first 12 months. The 3–6 month window "
     "(91–180 days) is the single largest risk bucket — 34 people lost, 31% of total. "
     "We are losing people before they become productive."),
    ("2",
     "ALMOST ALL UNDESIRED TALENT",
     "85% of exits (92 of 108) were people we wanted to keep. Only 15% were planned "
     "or desired exits. We are not in control of who leaves."),
    ("3",
     "PERSONAL & HEALTH = COVER STORY",
     "65% cite Personal or Health reasons — the two categories most often used when "
     "employees don't want to reveal the real reason (workload, manager, culture). "
     "This needs structured exit interviews and pulse surveys to uncover root causes."),
]

for i, (num, title, body) in enumerate(problems):
    y = 1.65 + i * 1.75
    add_rect(slide, 0.3, y, 12.7, 1.6, fill=RGBColor(0x14, 0x08, 0x08))
    add_rect(slide, 0.3, y, 0.8, 1.6, fill=RED_ALERT)
    add_textbox(slide, 0.3, y, 0.8, 1.6, num,
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 1.25, y + 0.12, 11.5, 0.42,
                title, font_size=14, bold=True, color=GOLD)
    add_textbox(slide, 1.25, y + 0.55, 11.5, 0.95,
                body, font_size=11, color=WHITE, wrap=True)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 8 — WHAT IS GOING RIGHT
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "What Is Going Right",
               "Recognise progress — build on it")

positives = [
    ("Absconds Down 50%",
     "From 14 in H2 2025 to just 7 in Q1 2026. Policy tightening and early-warning "
     "systems are clearly having an effect. Continue current practices.",
     "14  →  7"),
    ("Desired Exits More Than Doubled",
     "From 5 (9%) in 2025 to 11 (22%) in 2026. Performance management processes "
     "are being applied more effectively. The org is taking back control of exits.",
     "5  →  11"),
    ("Undesired Exit % Improving",
     "Dropped from 91% in H2 2025 to 78% in Q1 2026. Still high in absolute terms, "
     "but the directional trend is positive. Continue focused efforts.",
     "91%  →  78%"),
]

for i, (title, body, metric) in enumerate(positives):
    y = 1.65 + i * 1.75
    add_rect(slide, 0.3, y, 12.7, 1.6, fill=RGBColor(0x06, 0x18, 0x10))
    add_rect(slide, 0.3, y, 0.18, 1.6, fill=GREEN_OK)
    add_textbox(slide, 0.65, y + 0.1, 8.8, 0.42,
                title, font_size=14, bold=True, color=GREEN_OK)
    add_textbox(slide, 0.65, y + 0.55, 8.8, 0.9,
                body, font_size=11, color=WHITE, wrap=True)
    add_rect(slide, 9.6, y + 0.35, 3.2, 0.9, fill=RGBColor(0x10, 0x30, 0x20))
    add_textbox(slide, 9.6, y + 0.35, 3.2, 0.9,
                metric, font_size=18, bold=True, color=GREEN_OK,
                align=PP_ALIGN.CENTER)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 9 — ACTIONS REQUIRED FROM LEADERSHIP
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_header_bar(slide, "3 Actions Required from Leadership",
               "Concrete decisions needed — not analysis")

actions = [
    ("APPROVE",
     "90-Day Onboarding Reinforcement Program",
     "Target the 3–6 month tenure cohort immediately. Structured programme with clear "
     "milestones, buddy assignment, and manager commitment required.",
     ACCENT_BLUE),
    ("MANDATE",
     "Skip-Level Conversations for All Band 1 Employees",
     "At 60, 90, and 180-day marks. This creates a direct feedback channel bypassing "
     "line managers, enabling early identification of flight-risk employees.",
     GOLD),
    ("COMMISSION",
     "Workload & Wellness Audit — Asset Living F&A Ops Team",
     "65% Personal + Health exits is a warning signal. Commission an independent audit "
     "to surface real workload, WLB, and culture issues. Do not wait for more exits.",
     RED_ALERT),
]

for i, (action_word, title, body, col) in enumerate(actions):
    y = 1.65 + i * 1.75
    add_rect(slide, 0.3, y, 12.7, 1.6, fill=RGBColor(0x08, 0x16, 0x2E))
    add_rect(slide, 0.3, y, 1.7, 1.6, fill=col)
    add_textbox(slide, 0.3, y, 1.7, 1.6, action_word,
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, 2.15, y + 0.1, 10.65, 0.42,
                title, font_size=14, bold=True, color=col)
    add_textbox(slide, 2.15, y + 0.55, 10.65, 0.95,
                body, font_size=11, color=WHITE, wrap=True)

add_footer(slide)


# ══════════════════════════════════════════════════════════════════
#  SLIDE 10 — SUMMARY / CLOSING
# ══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_slide_bg(slide)
add_rect(slide, 0, 0, 0.6, 7.5, fill=ACCENT_BLUE)
add_rect(slide, 0.6, 0, 0.12, 7.5, fill=GOLD)

add_textbox(slide, 1.0, 0.7, 11.5, 0.6,
            "SUMMARY — Asset Living F&A Attrition Report", font_size=26,
            bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 1.0, 1.3, 11.5, 0.38,
            "July 2025 – March 2026  |  108 Total Exits",
            font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF),
            align=PP_ALIGN.CENTER, italic=True)

summary_items = [
    (RED_ALERT,  "2026 monthly rate is 17/mo — 79% higher than H2 2025 (9.5/mo). Pace must be reversed."),
    (RED_ALERT,  "57% of exits occur in the first 12 months. The 91–180 day window is the highest risk."),
    (RED_ALERT,  "85% of exits are undesired — we are not in control of who leaves."),
    (RED_ALERT,  "65% cite Personal/Health — likely a proxy for workload, culture, or manager issues."),
    (GREEN_OK,   "Absconds halved (14 → 7). Policy measures are working. Sustain the effort."),
    (GREEN_OK,   "Desired exits improving (5 → 11). Performance management is being used effectively."),
    (GOLD,       "Immediate action: 90-day reinforcement programme  |  Skip-level check-ins  |  Wellness audit"),
]

for i, (col, text) in enumerate(summary_items):
    y = 2.05 + i * 0.65
    add_rect(slide, 1.0, y + 0.12, 0.25, 0.28, fill=col)
    add_textbox(slide, 1.4, y, 11.3, 0.6, text, font_size=12,
                color=WHITE if col != GREEN_OK else RGBColor(0xAA, 0xFF, 0xCC),
                wrap=True)

add_footer(slide)


# ── Save ────────────────────────────────────────────────────────
out = "/Applications/XAMPP/xamppfiles/htdocs/minerva/tools/Attrition_Analysis_Report_Asset_Living.pptx"
prs.save(out)
print("Saved:", out)
